# backend/app/services/export_ppt.py
"""
Rombak total (ganti gaya lama sepenuhnya): palet hijau/emas PT Petrokimia Gresik, font
Bookman Old Style (judul) + Calibri (body), TANPA bullet titik (badge lingkaran nomor/
huruf), TANPA em dash (sanitize_text), TANPA garis aksen/bar dekoratif di judul/kartu/
footer (kecuali ornamen lengkung emas di cover & penutup, itu bukan "garis aksen" yang
dilarang — cuma flourish sudut satu kali). Chart NATIVE python-pptx (bukan gambar PNG dari
Plotly/Kaleido — chart_generator.py TIDAK disentuh/dipakai lagi di sini, statistik chart
diambil LANGSUNG dari compute_statistics()).

Jumlah & kehadiran slide FLEKSIBEL mengikuti data yang benar-benar tersedia (skip aman
kalau kolom terkait tak terdeteksi) — bukan struktur 12-slide yang kaku.

Konten narasi HANYA dari 6 key wajib lama (executive_summary, dst) + key_findings opsional
— SENGAJA TIDAK memakai ai_summary["sections"] (PART A) supaya render ini otomatis
backward & forward compatible tanpa menyentuh fitur section dinamis sama sekali.

Layout bervariasi ANTAR LAPORAN (posisi panel, jumlah kolom grid, gaya cover/chart/kartu,
sudut ornamen) — dipilih SEKALI per laporan lewat pick_visual_style() (report_render_logic.py)
tepat saat analisis AI berhasil, disimpan ke report.visual_style, dibaca di sini via
get_visual_style() (BUKAN di-random di sini lagi) supaya preview web & file yang diunduh
SELALU menampilkan bentuk yang identik utk laporan yang sama, tapi tetap dalam identitas
visual (palet/font/makna warna) yang sama.
"""
import math
from dataclasses import dataclass

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

from app.models.report import Report
from app.services.report_render_logic import build_report_blocks, is_english, find_logo_path, get_visual_style, resolve_theme_color

# ============================================================================
# Palet & font — persis sesuai brief, dipakai di SETIAP elemen (termasuk chart/tabel)
# ============================================================================
GREEN_MAIN = RGBColor(0x1B, 0x5E, 0x3C)
GREEN_BG = RGBColor(0x0E, 0x3B, 0x26)
GREEN_CHART = RGBColor(0x2F, 0x7A, 0x52)
GOLD_MAIN = RGBColor(0xC9, 0xA2, 0x27)
GOLD_LIGHT = RGBColor(0xE7, 0xC7, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
IVORY = RGBColor(0xF5, 0xF7, 0xF2)
TEXT_DARK = RGBColor(0x16, 0x24, 0x1C)
GRAY_TEXT = RGBColor(0x5C, 0x6B, 0x62)
RED_CRIT = RGBColor(0xB2, 0x3A, 0x2E)
RED_CRIT_BG = RGBColor(0xF8, 0xE2, 0xDE)
PANEL_BORDER = RGBColor(0xE2, 0xE5, 0xDE)
TITLE_FONT = "Bookman Old Style"
BODY_FONT = "Calibri"

# ── TEMA WARNA (report.theme_color) ─────────────────────────────────────────
# Sama persis dengan export_pdf.py — GREEN_MAIN/BG/CHART & GOLD_MAIN/LIGHT di atas TETAP
# dipakai langsung oleh SEVERITY_COLOR di bawah (warna severity TIDAK ikut tema apa pun).
# THEME_PALETTES murni untuk elemen BRAND/struktural (cover, kicker, badge, border panel,
# header tabel, chart "bar" utama) — nilai HEX identik dengan export_pdf.py/reportTheme.ts.
NAVY_MAIN = RGBColor(0x1E, 0x3A, 0x5F)
NAVY_BG = RGBColor(0x0F, 0x17, 0x2A)
NAVY_CHART = RGBColor(0x3B, 0x6E, 0xA5)
DARK_MAIN = RGBColor(0x1F, 0x29, 0x37)
DARK_BG = RGBColor(0x11, 0x18, 0x27)
DARK_CHART = RGBColor(0x3F, 0x4B, 0x5C)
GOLD_BRONZE_MAIN = RGBColor(0x8A, 0x6A, 0x16)
GOLD_BRONZE_BG = RGBColor(0x4A, 0x39, 0x08)
GOLD_CREAM_LIGHT = RGBColor(0xF3, 0xE3, 0xAE)
GOLD_CREAM_SOFT = RGBColor(0xFB, 0xF3, 0xDC)

THEME_PALETTES: dict[str, dict[str, RGBColor]] = {
    "green": {"main": GREEN_MAIN, "bg": GREEN_BG, "chart": GREEN_CHART, "light": GOLD_MAIN, "soft": GOLD_LIGHT},
    "navy": {"main": NAVY_MAIN, "bg": NAVY_BG, "chart": NAVY_CHART, "light": GOLD_MAIN, "soft": GOLD_LIGHT},
    "dark": {"main": DARK_MAIN, "bg": DARK_BG, "chart": DARK_CHART, "light": GOLD_MAIN, "soft": GOLD_LIGHT},
    "gold": {"main": GOLD_BRONZE_MAIN, "bg": GOLD_BRONZE_BG, "chart": GOLD_MAIN, "light": GOLD_CREAM_LIGHT, "soft": GOLD_CREAM_SOFT},
}

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
MARGIN_X = Inches(0.5)
CONTENT_W = SLIDE_W - MARGIN_X * 2

# Warna ramp kategori (dipakai legend panel "% Proporsi Kategori") — nuansa hijau/emas.
CATEGORY_COLOR_RAMP = [GREEN_MAIN, GREEN_CHART, GOLD_MAIN, GOLD_LIGHT, GRAY_TEXT]

# Warna per-level severity (bar chart & badge) — Critical selalu merah, tidak pernah "aman".
SEVERITY_COLOR = {
    "critical": RED_CRIT,
    "high": GOLD_MAIN,
    "medium": GREEN_MAIN,
    "low": GREEN_CHART,
    "informational": GRAY_TEXT,
}
def _resolve_logo_path() -> str | None:
    return find_logo_path()


# ============================================================================
# Helper visual dasar
# ============================================================================
def _set_font(para_or_run, name=BODY_FONT, size=None, bold=None, italic=None, color=None):
    """Terima paragraph ATAU run. Kalau paragraph SUDAH punya run (biasanya karena
    `.text` sudah di-set, jadi ada 1 run implisit), font diterapkan LANGSUNG ke tiap run
    -- bukan cuma default level-paragraf (defRPr) -- supaya font benar-benar tertanam &
    konsisten terbaca (mis. oleh alat verifikasi atau software office lain)."""
    targets = para_or_run.runs if hasattr(para_or_run, "runs") and para_or_run.runs else [para_or_run]
    for t in targets:
        f = t.font
        f.name = name
        if size is not None:
            f.size = size
        if bold is not None:
            f.bold = bold
        if italic is not None:
            f.italic = italic
        if color is not None:
            f.color.rgb = color


def _estimate_wrapped_height_in(text: str, font_pt: float, box_width_in: float) -> float:
    """Estimasi kasar tinggi (inci) yang dipakai `text` kalau word-wrap dalam box selebar
    `box_width_in` pada ukuran `font_pt` — python-pptx bukan mesin render sungguhan (tidak
    tahu lebar karakter asli tiap font), jadi ini estimasi berbasis rata-rata lebar karakter
    huruf tebal (0.6em). Dipakai supaya elemen SETELAH judul (subtitle, info) tidak pernah
    ketimpa kalau judulnya panjang & wrap ke banyak baris — laporan bisa dari domain apa saja
    (SOC, keuangan, KPI, dll) dengan panjang judul yang jauh berbeda-beda, jadi TIDAK BOLEH
    diasumsikan selalu pendek/selalu 1 baris."""
    if not text:
        return font_pt * 1.25 / 72
    avg_char_width_in = (font_pt * 0.6) / 72
    chars_per_line = max(int(box_width_in / avg_char_width_in), 1)
    line_count = max(1, -(-len(text) // chars_per_line))  # ceil division
    return line_count * (font_pt * 1.25) / 72


def _no_shadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def _modest_corner(shape, frac=0.06):
    """PowerPoint memberi ROUNDED_RECTANGLE radius sudut default besar (~1/6 dari sisi
    terpendek) — utk shape yang tidak terlalu tinggi/lebar, itu bikin lengkungan sudut
    memakan inset konten di dekatnya (badge/teks pojok kiri-atas terlihat menempel/tertimpa
    lengkungan, dilaporkan user lewat tangkapan layar). Dipanggil setelah tiap ROUNDED_
    RECTANGLE dibuat (kecuali yang SENGAJA dibuat pil penuh via adjustments[0]=0.5)."""
    try:
        shape.adjustments[0] = frac
    except Exception:
        pass


def _send_to_back(slide, shape):
    sp = shape._element
    spTree = slide.shapes._spTree
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_dark_bg(slide, theme: dict | None = None):
    t = theme or THEME_PALETTES["green"]
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = t["bg"]
    rect.line.fill.background()
    _no_shadow(rect)
    _send_to_back(slide, rect)
    return rect


def _fill_rect_bg(slide, x, y, w, h, color):
    """Persegi latar penuh (bukan slide utuh) — dipakai cover/penutup varian split-warna
    (`cover_style="split"`) utk 2 kolom warna solid berdampingan, analog `add_dark_bg` tapi
    utk area sebagian slide."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    _no_shadow(rect)
    _send_to_back(slide, rect)
    return rect


def add_corner_flourish(slide, corner: str = "bottom_right", area_x=0, area_y=0, area_w=None, area_h=None, theme: dict | None = None):
    """Ornamen lengkung emas tipis (beberapa lingkaran konsentris tanpa isi, diposisikan
    menjorok keluar sudut) — dipakai HANYA di cover & penutup, mendekati motif referensi.

    `area_x/y/w/h` (opsional, default = seluruh slide) membatasi sudut mana yang dipakai
    acuan — dibutuhkan utk cover/penutup varian split-warna (`cover_style="split"`), di mana
    flourish HARUS tetap berada di dalam kolom hijau saja (bukan di sudut fisik slide penuh,
    yang sebagian jatuh di kolom emas kalau corner="bottom_left")."""
    t = theme or THEME_PALETTES["green"]
    area_w = SLIDE_W if area_w is None else area_w
    area_h = SLIDE_H if area_h is None else area_h
    # Radius & inset dikecilkan dari versi awal (dulu sampai 3.25in, base offset 0.8/0.6in) —
    # keluhan nyata dari pengguna: lingkaran terluar menjorok cukup jauh sampai menembus area
    # teks "Diskusi dan pertanyaan dipersilakan" di penutup, mengganggu keterbacaan. Ornamen
    # tetap ada (identitas visual), cuma jangkauannya dipersempit supaya tetap di pojok saja.
    if corner == "bottom_left":
        base_x, base_y = area_x - Inches(0.8), area_y + area_h - Inches(1.0)
    elif corner == "top_right":
        base_x, base_y = area_x + area_w - Inches(1.0), area_y - Inches(0.8)
    else:
        base_x, base_y = area_x + area_w - Inches(1.0), area_y + area_h - Inches(1.0)

    for r in (Inches(1.3), Inches(1.75), Inches(2.2), Inches(2.65)):
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, base_x - r, base_y - r, r * 2, r * 2)
        oval.fill.background()
        oval.line.color.rgb = t["light"]
        oval.line.width = Pt(0.75)
        _no_shadow(oval)


def add_kicker(slide, text: str, color=GREEN_MAIN, x=MARGIN_X, y=Inches(0.35), w=Inches(9)):
    box = slide.shapes.add_textbox(x, y, w, Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = text.upper()
    _set_font(p, BODY_FONT, Pt(11), bold=True, color=color)
    return box


def add_title(slide, text: str, color=TEXT_DARK, x=MARGIN_X, y=Inches(0.68), w=Inches(11.5), size=Pt(30)):
    """Return: posisi Y (inci) TEPAT DI BAWAH judul ini — WAJIB dipakai pemanggil utk
    memposisikan elemen berikutnya (bukan konstanta tetap seperti Inches(1.45) dst).

    BUG BESAR YANG DIPERBAIKI: box judul SEBELUMNYA selalu Inches(0.75) tinggi tetap,
    diam-diam mengasumsikan judul selalu muat 1 baris. Kalau block["title"]/["heading"]
    panjang (umum utk data non-SOC — nama bulan+tahun, kalimat AI, dst) dan wrap ke 2 baris,
    baris kedua itu VISUAL MELUBER ke luar box (PowerPoint textbox TIDAK auto-reflow shape
    lain di baliknya) dan menabrak elemen berikutnya yang posisinya konstanta tetap (terlihat
    nyata di slide Ringkasan Eksekutif: baris ke-2 judul bertabrakan dgn kartu KPI di
    bawahnya). Tinggi box SEKARANG dihitung dari estimasi wrap sungguhan (_estimate_wrapped_
    height_in, sudah dipakai di file ini utk hal serupa), dan Y bawahnya DIKEMBALIKAN supaya
    tiap slide bisa menggeser elemen berikutnya sesuai tinggi judul yang SEBENARNYA."""
    box = slide.shapes.add_textbox(x, y, w, Inches(0.75))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    _set_font(p, TITLE_FONT, size, bold=True, color=color)
    text_h_in = _estimate_wrapped_height_in(text, size.pt, Emu(w).inches)
    box_h_in = max(0.75, text_h_in + 0.08)
    box.height = Inches(box_h_in)
    return Emu(y).inches + box_h_in


def add_footer(slide, page_num: int, total_pages: int):
    box = slide.shapes.add_textbox(MARGIN_X, SLIDE_H - Inches(0.42), CONTENT_W, Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = f"{page_num:02d} / {total_pages:02d}"
    p.alignment = PP_ALIGN.RIGHT
    _set_font(p, BODY_FONT, Pt(9), color=GRAY_TEXT)


def add_logo(slide, logo_path, x=None, y=Inches(0.25), width=Inches(1.7)):
    if not logo_path:
        return
    x = x if x is not None else (SLIDE_W - width - Inches(0.35))
    try:
        slide.shapes.add_picture(logo_path, x, y, width=width)
    except Exception:
        pass


def add_badge_circle(slide, x, y, diameter, text, badge_color, text_color=WHITE, font_size=Pt(14)):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, diameter, diameter)
    circle.fill.solid()
    circle.fill.fore_color.rgb = badge_color
    circle.line.fill.background()
    _no_shadow(circle)
    tf = circle.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = str(text)
    p.alignment = PP_ALIGN.CENTER
    _set_font(p, BODY_FONT, font_size, bold=True, color=text_color)
    return circle


def add_badge_row(slide, x, y, w, number_text, title_text, detail_text, badge_color,
                   badge_d=Inches(0.42), on_dark=False, title_size=Pt(15), detail_size=Pt(12)):
    """Return: perkiraan tinggi konten (title+detail, dalam inci) — dipakai add_badge_list
    untuk menggeser baris berikutnya sesuai tinggi SEBENARNYA, bukan jarak tetap."""
    add_badge_circle(slide, x, y, badge_d, number_text, badge_color, font_size=Pt(14))
    text_x = x + badge_d + Inches(0.2)
    text_w = w - badge_d - Inches(0.2)
    box = slide.shapes.add_textbox(text_x, y - Inches(0.03), text_w, Inches(0.95))
    tf = box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = title_text
    _set_font(p1, BODY_FONT, title_size, bold=True, color=(WHITE if on_dark else TEXT_DARK))
    text_w_in = Emu(text_w).inches
    content_height_in = _estimate_wrapped_height_in(title_text, title_size.pt, text_w_in)
    if detail_text:
        p2 = tf.add_paragraph()
        p2.text = detail_text
        _set_font(p2, BODY_FONT, detail_size, color=(GOLD_LIGHT if on_dark else GRAY_TEXT))
        p2.space_before = Pt(2)
        content_height_in += _estimate_wrapped_height_in(detail_text, detail_size.pt, text_w_in) + 2 / 72
    return content_height_in


def _estimate_badge_row_height_in(w_in, badge_d_in, title_text, detail_text, title_pt, detail_pt):
    """Estimasi tinggi 1 add_badge_row TANPA menggambar apapun (fungsi murni) — dipakai
    add_badge_list utk pre-pass menghitung total tinggi SEMUA item SEBELUM mulai render."""
    text_w_in = w_in - badge_d_in - 0.2
    h = _estimate_wrapped_height_in(title_text, title_pt, text_w_in)
    if detail_text:
        h += _estimate_wrapped_height_in(detail_text, detail_pt, text_w_in) + 2 / 72
    return h


def add_badge_list(slide, x, y, w, items, badge_color=GREEN_MAIN, row_h=Inches(0.95), on_dark=False, max_y=None, cols=1):
    """items: list of (number_or_letter, title, detail). badge_color: RGBColor tetap ATAU
    fungsi(idx, item)->RGBColor supaya bisa mem-variasi warna (mis. merah utk 1 item kritis).
    `row_h` dipakai sebagai jarak MINIMUM antar baris saja — kalau title/detail dari AI
    kebetulan lebih panjang dan wrap ke banyak baris, jaraknya digeser sesuai perkiraan tinggi
    sebenarnya (dulu selalu jarak tetap, baris berikutnya bisa menimpa baris ini).

    `cols` (default 1, backward compat) — BUG NYATA YANG DIPERBAIKI (dilaporkan user): dipanggil
    dengan `w`=CONTENT_W penuh (mis. Temuan Utama, sampai 6 item) menghasilkan baris selebar
    HAMPIR SELURUH slide (~11.8in dari 13.3in) — capek dibaca, apalagi digabung skala-turun font
    di bawah. `cols=2` membelah item jadi grid 2 kolom (lebar per baris otomatis separuh, jauh
    di bawah ambang ~75 karakter/baris yang wajar dibaca), DAN memangkas jumlah BARIS jadi
    separuhnya — dua-duanya mengurangi kebutuhan skala-turun font drastis dibanding 1 kolom
    penuh utk jumlah item yang sama.

    `max_y` (opsional, disarankan selalu diisi) = batas bawah yang TIDAK BOLEH dilewati (mis.
    SLIDE_H dikurangi margin footer). BUG YANG DIPERBAIKI: sebelum ini tidak ada pengecekan
    apapun terhadap tinggi slide — laporan dengan item lebih banyak (mis. 6 Temuan Utama)
    membuat item ke-5/6 punya posisi Y melebihi slide_height, jadi tidak terlihat sama sekali
    saat dibuka di PowerPoint (bukan error, cuma diam-diam hilang). Sekarang total tinggi
    SEMUA BARIS (bukan item — 1 baris bisa berisi `cols` item berjajar) diperkirakan DULU
    (pre-pass, tanpa menggambar) sebelum mulai render — kalau ternyata bakal melebihi `max_y`,
    seluruh baris (badge, judul, detail, jarak antar baris) dikecilkan skalanya secara
    proporsional supaya SEMUA item pasti muat di dalam slide."""
    if not items:
        return y
    row_h_in = row_h.inches
    col_gap_in = 0.4
    col_w_in = (Emu(w).inches - col_gap_in * (cols - 1)) / cols
    col_w = Inches(col_w_in)
    default_badge_d_in = 0.42
    default_title_pt, default_detail_pt = 15, 12
    n_rows = math.ceil(len(items) / cols)

    scale = 1.0
    if max_y is not None:
        available_in = Emu(max_y - y).inches
        row_heights_in = [
            max(
                row_h_in,
                max(
                    _estimate_badge_row_height_in(col_w_in, default_badge_d_in, it[1], it[2], default_title_pt, default_detail_pt) + 0.15
                    for it in items[r * cols: r * cols + cols]
                ),
            )
            for r in range(n_rows)
        ]
        est_total_in = sum(row_heights_in)
        if available_in > 0 and est_total_in > available_in:
            scale = max(available_in / est_total_in, 0.55)

    title_size = Pt(default_title_pt * scale)
    detail_size = Pt(default_detail_pt * scale)
    badge_d = Inches(default_badge_d_in * scale)
    row_h_scaled_in = row_h_in * scale
    row_gap_in = 0.15 * scale

    cur_y = y
    for r in range(n_rows):
        row_start = r * cols
        row_items = items[row_start:row_start + cols]
        max_content_h_in = 0.0
        for c, item in enumerate(row_items):
            idx = row_start + c
            color = badge_color(idx, item) if callable(badge_color) else badge_color
            cx = x + c * (col_w + Inches(col_gap_in))
            content_height_in = add_badge_row(slide, cx, cur_y, col_w, item[0], item[1], item[2], color,
                                               badge_d=badge_d, on_dark=on_dark,
                                               title_size=title_size, detail_size=detail_size)
            max_content_h_in = max(max_content_h_in, content_height_in)
        cur_y += Inches(max(row_h_scaled_in, max_content_h_in + row_gap_in))
    return cur_y


def add_stat_card_grid(slide, x, y, w, h, items, cols=3, dark=True, theme: dict | None = None):
    """items: list of (value_str, label_str)."""
    t = theme or THEME_PALETTES["green"]
    if not items:
        return y
    rows = math.ceil(len(items) / cols)
    gap = Inches(0.2)
    # card_h dibatasi maksimum — TANPA batas ini, laporan dengan sedikit stat item (mis. cuma 2
    # kartu KPI utk domain yang tidak punya konsep severity/status) membuat 1 baris itu meregang
    # mengisi SELURUH `h` yang dialokasikan pemanggil (mis. 4.3in), kartu raksasa nyaris kosong
    # utk cuma 1 angka + label pendek. Batasnya SEKARANG bertingkat sesuai jumlah baris (dulu
    # konstan 1.6in utk semua kasus) — laporan dengan SEDIKIT baris (ruang per kartu otomatis
    # lebih longgar) dapat kartu lebih tinggi + font lebih besar supaya ruang terasa terisi
    # proporsional, BUKAN cuma dibuat kecil rata utk menghindari kartu raksasa (keluhan nyata
    # dari pengguna: slide "Ringkasan Eksekutif" isinya cuma sepertiga atas, sisanya kosong).
    max_card_h = {1: Inches(2.6), 2: Inches(2.0)}.get(rows, Inches(1.6))
    natural_card_h = (h - gap * (rows - 1)) / rows
    card_h = min(natural_card_h, max_card_h)
    card_h_in = Emu(card_h).inches
    value_pt = 40 if card_h_in >= 2.2 else 36 if card_h_in >= 1.8 else 32
    label_pt = 13 if card_h_in >= 2.2 else 12 if card_h_in >= 1.8 else 11.5
    for idx, (value, label) in enumerate(items):
        r = idx // cols
        row_start = r * cols
        # Lebar kartu dihitung dari JUMLAH ITEM DI BARIS INI, bukan `cols` tetap — supaya
        # baris terakhir yang isinya lebih sedikit dari `cols` (mis. cuma 1 kartu tersisa)
        # melebar mengisi ruang, bukan tampil sempit dengan sisa ruang kosong di sampingnya.
        row_items_count = min(cols, len(items) - row_start)
        c = idx - row_start
        card_w = (w - gap * (row_items_count - 1)) / row_items_count
        cx = x + c * (card_w + gap)
        cy = y + r * (card_h + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = t["main"] if dark else IVORY
        card.line.color.rgb = t["light"]
        card.line.width = Pt(0.75)
        _no_shadow(card)
        _modest_corner(card)
        tf = card.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = str(value)
        p1.alignment = PP_ALIGN.CENTER
        _set_font(p1, TITLE_FONT, Pt(value_pt), bold=True, color=t["light"] if dark else t["main"])
        p2 = tf.add_paragraph()
        p2.text = label
        p2.alignment = PP_ALIGN.CENTER
        _set_font(p2, BODY_FONT, Pt(label_pt), color=WHITE if dark else TEXT_DARK)
        p2.space_before = Pt(6)
    return y + rows * card_h + (rows - 1) * gap


def add_ivory_panel(slide, x, y, w, h, icon_text, title_text, rows, mode="kv", footnote=None, theme: dict | None = None):
    """mode="kv": rows = [(label, value), ...]. mode="legend": rows = [(color, label, pct), ...]."""
    t = theme or THEME_PALETTES["green"]
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    panel.fill.solid()
    panel.fill.fore_color.rgb = IVORY
    panel.line.color.rgb = PANEL_BORDER
    panel.line.width = Pt(0.75)
    _no_shadow(panel)
    _modest_corner(panel)  # lihat docstring _modest_corner — bug badge menempel border

    pad = Inches(0.3)
    inner_x = x + pad
    inner_w = w - pad * 2
    cur_y = y + Inches(0.26)

    add_badge_circle(slide, inner_x, cur_y, Inches(0.32), icon_text, t["light"], font_size=Pt(11))
    title_box = slide.shapes.add_textbox(inner_x + Inches(0.45), cur_y + Inches(0.02), inner_w - Inches(0.45), Inches(0.32))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title_text.upper()
    _set_font(tp, BODY_FONT, Pt(11.5), bold=True, color=t["main"])
    cur_y += Inches(0.55)

    # BUG YANG DIPERBAIKI: dulu tidak ada pengecekan terhadap `h` (tinggi panel tetap dari
    # pemanggil) — panel dengan banyak baris (mis. legend 6 kategori + footnote) berisiko
    # baris terakhirnya meluber ke luar panel. Total tinggi SEMUA baris diperkirakan DULU
    # (pre-pass), lalu kalau bakal melebihi sisa ruang panel, jarak antar baris & ukuran
    # font dikecilkan proporsional supaya semua baris tetap muat di dalam panel.
    available_in = Emu(h).inches - Emu(cur_y - y).inches - 0.2 - (0.3 if footnote else 0)
    scale = 1.0
    inner_w_in = Emu(inner_w).inches
    if mode == "kv":
        est_total_in = sum(max(0.62, 0.22 + _estimate_wrapped_height_in(str(v), 11.5, inner_w_in) + 0.16) for _, v in rows)
    else:
        label_w_in_est = Emu(inner_w - Inches(1.0)).inches
        est_total_in = sum(max(0.34, _estimate_wrapped_height_in(label, 11, label_w_in_est) + 0.06) for _, label, _ in rows)
    if available_in > 0 and est_total_in > available_in:
        scale = max(available_in / est_total_in, 0.55)

    if mode == "kv":
        # Jarak antar baris menyesuaikan tinggi VALUE sebenarnya (bukan 0.62in tetap) — value
        # bisa berisi teks bebas dengan panjang tidak menentu (mis. nama file yang diunggah
        # pengguna), yang tanpa ini berisiko menimpa baris berikutnya kalau kebetulan panjang.
        label_pt, value_pt = 11.5 * scale, 11.5 * scale
        for label, value in rows:
            # BUG YANG DIPERBAIKI (dilaporkan user): alignment tidak pernah di-set eksplisit
            # di sini — label & value jadi berisiko dirender beda pola (mis. label rata tengah,
            # value rata kiri) tergantung default theme, membuat teks tampak zigzag tidak rapi.
            # Keduanya SEKARANG eksplisit rata kiri, pola yang sama, konsisten satu sama lain.
            lbl_box = slide.shapes.add_textbox(inner_x, cur_y, inner_w, Inches(0.24))
            lp = lbl_box.text_frame.paragraphs[0]
            lp.text = label
            lp.alignment = PP_ALIGN.LEFT
            _set_font(lp, BODY_FONT, Pt(label_pt), bold=True, color=t["main"])
            val_box = slide.shapes.add_textbox(inner_x, cur_y + Inches(0.22 * scale), inner_w, Inches(0.42))
            vtf = val_box.text_frame
            vtf.word_wrap = True
            vp = vtf.paragraphs[0]
            vp.text = str(value)
            vp.alignment = PP_ALIGN.LEFT
            _set_font(vp, BODY_FONT, Pt(value_pt), color=GRAY_TEXT)
            value_height_in = _estimate_wrapped_height_in(str(value), value_pt, inner_w_in)
            cur_y += Inches(max(0.62 * scale, (0.22 + value_height_in + 0.16) * scale))
    else:
        label_w_in = Emu(inner_w - Inches(1.0)).inches
        label_pt = 11 * scale
        for color, label, pct in rows:
            sw = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inner_x, cur_y + Inches(0.03), Inches(0.14), Inches(0.14))
            sw.fill.solid()
            sw.fill.fore_color.rgb = color
            sw.line.fill.background()
            _no_shadow(sw)
            lbl_box = slide.shapes.add_textbox(inner_x + Inches(0.22), cur_y, inner_w - Inches(1.0), Inches(0.28))
            lp = lbl_box.text_frame.paragraphs[0]
            lp.text = label
            _set_font(lp, BODY_FONT, Pt(label_pt), color=TEXT_DARK)
            val_box = slide.shapes.add_textbox(inner_x + inner_w - Inches(0.8), cur_y, Inches(0.8), Inches(0.28))
            vp = val_box.text_frame.paragraphs[0]
            vp.text = pct
            vp.alignment = PP_ALIGN.RIGHT
            _set_font(vp, BODY_FONT, Pt(label_pt), bold=True, color=t["main"])
            label_height_in = _estimate_wrapped_height_in(label, label_pt, label_w_in)
            cur_y += Inches(max(0.34 * scale, (label_height_in + 0.06) * scale))

    if footnote:
        cur_y += Inches(0.12)
        divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inner_x, cur_y, inner_w, Pt(0.75))
        divider.fill.solid()
        divider.fill.fore_color.rgb = PANEL_BORDER
        divider.line.fill.background()
        _no_shadow(divider)
        cur_y += Inches(0.16)
        fn_box = slide.shapes.add_textbox(inner_x, cur_y, inner_w, max(y + h - cur_y - Inches(0.1), Inches(0.3)))
        ftf = fn_box.text_frame
        ftf.word_wrap = True
        fp = ftf.paragraphs[0]
        fp.text = footnote
        _set_font(fp, BODY_FONT, Pt(9.5), italic=True, color=GRAY_TEXT)
    return panel


def add_dark_panel(slide, x, y, w, h, theme: dict | None = None):
    t = theme or THEME_PALETTES["green"]
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    panel.fill.solid()
    panel.fill.fore_color.rgb = t["bg"]
    panel.line.color.rgb = t["light"]
    panel.line.width = Pt(1)
    _no_shadow(panel)
    _modest_corner(panel)
    return panel


def add_critical_highlight_panel(slide, x, y, w, h, pct_text, sub_text, detail_text=None, theme: dict | None = None):
    t = theme or THEME_PALETTES["green"]
    add_dark_panel(slide, x, y, w, h, theme=t)
    pad = Inches(0.3)
    big_box = slide.shapes.add_textbox(x + pad, y + Inches(0.35), w - pad * 2, Inches(1.0))
    bp = big_box.text_frame.paragraphs[0]
    bp.text = pct_text
    bp.alignment = PP_ALIGN.CENTER
    _set_font(bp, TITLE_FONT, Pt(42), bold=True, color=t["light"])

    sub_top_in = 1.35
    sub_box = slide.shapes.add_textbox(x + pad, y + Inches(sub_top_in), w - pad * 2, Inches(0.7))
    stf = sub_box.text_frame
    stf.word_wrap = True
    sp = stf.paragraphs[0]
    sp.text = sub_text
    sp.alignment = PP_ALIGN.CENTER
    _set_font(sp, BODY_FONT, Pt(12.5), color=WHITE)

    if detail_text:
        # Posisi detail_text dihitung dari perkiraan tinggi sub_text sebenarnya (bukan y+2.15
        # tetap) — sub_text panjangnya tidak menentu (kalimat dari data), jadi tanpa ini
        # detail_text (daftar kategori insiden Critical, bisa sampai 6 nama) berisiko tertimpa.
        text_w_in = Emu(w - pad * 2).inches
        sub_height_in = _estimate_wrapped_height_in(sub_text, 12.5, text_w_in)
        detail_top_in = max(sub_top_in + sub_height_in + 0.15, 2.15)
        det_box = slide.shapes.add_textbox(x + pad, y + Inches(detail_top_in), w - pad * 2, h - Inches(detail_top_in) - Inches(0.25))
        dtf = det_box.text_frame
        dtf.word_wrap = True
        dp = dtf.paragraphs[0]
        dp.text = detail_text
        _set_font(dp, BODY_FONT, Pt(10.5), color=t["soft"])


def add_priority_panel(slide, x, y, w, h, title_text, items, theme: dict | None = None):
    """items: list of (letter, text)."""
    t = theme or THEME_PALETTES["green"]
    add_dark_panel(slide, x, y, w, h, theme=t)
    pad = Inches(0.28)
    title_box = slide.shapes.add_textbox(x + pad, y + Inches(0.22), w - pad * 2, Inches(0.32))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title_text.upper()
    _set_font(tp, BODY_FONT, Pt(11.5), bold=True, color=t["light"])
    if not items:
        return
    content_start_y_in = Emu(y).inches + 0.7
    available_in = Emu(h).inches - 0.7 - 0.2
    text_box_w_in = Emu(w - pad * 2 - Inches(0.5)).inches

    # BUG YANG DIPERBAIKI: tidak ada pengecekan sebelumnya terhadap tinggi panel `h` (fixed,
    # dari pemanggil) — item yang cukup banyak/panjang bisa membuat baris terakhir meluber ke
    # luar panel (bahkan ke luar slide). Sekarang total tinggi diperkirakan DULU (pre-pass),
    # dan kalau bakal melebihi `h`, badge+teks+jarak antar baris dikecilkan proporsional.
    scale = 1.0
    est_total_in = sum(max(0.62, _estimate_wrapped_height_in(text, 13, text_box_w_in) + 0.2) for _, text in items)
    if available_in > 0 and est_total_in > available_in:
        scale = max(available_in / est_total_in, 0.55)

    badge_d = Inches(0.36 * scale)
    font_pt = 13 * scale
    row_min_in = 0.62 * scale
    cur_y = Inches(content_start_y_in)
    for letter, text in items:
        add_badge_circle(slide, x + pad, cur_y, badge_d, letter, t["light"], font_size=Pt(13 * scale))
        box = slide.shapes.add_textbox(x + pad + Inches(0.5), cur_y + Inches(0.02), w - pad * 2 - Inches(0.5), Inches(0.55))
        btf = box.text_frame
        btf.word_wrap = True
        bp = btf.paragraphs[0]
        bp.text = text
        _set_font(bp, BODY_FONT, Pt(font_pt), color=WHITE)
        # Baris berikutnya digeser sesuai perkiraan tinggi teks yang SEBENARNYA (bukan jarak
        # tetap) — teks rekomendasi dari AI panjangnya tidak menentu, dan jarak tetap bikin
        # baris berikutnya menimpa baris ini kalau teksnya wrap lebih dari ~2 baris.
        text_height_in = _estimate_wrapped_height_in(text, font_pt, text_box_w_in)
        cur_y += Inches(max(row_min_in, text_height_in * scale + 0.2 * scale))


def add_ai_insight_strip(slide, x, y, w, text):
    """Kotak singkat "Insight AI" di bawah chart — menampilkan chart_captions dari AI, yang
    SEBELUMNYA dihasilkan AI (lihat prompts.py) tapi tidak pernah ditampilkan di PDF/PPT
    sama sekali. Cuma dipanggil kalau ai_caption benar-benar ada isinya."""
    box = slide.shapes.add_textbox(x, y, w, Inches(0.6))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"\U0001F4A1 {text}"
    _set_font(p, BODY_FONT, Pt(10), italic=True, color=GRAY_TEXT)
    return box


def add_pill_stat(slide, x, y, w, h, text, theme: dict | None = None):
    t = theme or THEME_PALETTES["green"]
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        pill.adjustments[0] = 0.5
    except Exception:
        pass
    pill.fill.solid()
    pill.fill.fore_color.rgb = t["main"]
    pill.line.color.rgb = t["light"]
    pill.line.width = Pt(1)
    _no_shadow(pill)
    tf = pill.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    _set_font(p, BODY_FONT, Pt(13), bold=True, color=t["light"])


def add_asset_card_row(slide, x, y, w, h, items, theme: dict | None = None):
    """items: list of (badge_num, title, stat_text, desc_text). `h` adalah tinggi ZONA yang
    dialokasikan pemanggil (bisa jauh lebih besar dari kebutuhan konten sebenarnya)."""
    t = theme or THEME_PALETTES["green"]
    n = len(items)
    if n == 0:
        return y
    gap = Inches(0.3)
    card_w = (w - gap * (n - 1)) / n
    pad = Inches(0.22)
    pad_in = 0.22
    text_w_in = Emu(card_w).inches - pad_in * 2

    # Tinggi kartu SEKARANG mengikuti kebutuhan konten sebenarnya (badge+judul+stat+deskripsi),
    # bukan selalu memenuhi penuh `h` yang dialokasikan pemanggil — keluhan nyata dari pengguna:
    # kartu vendor kosong ~2/3 bagian bawahnya kalau deskripsinya pendek/`h` yang dialokasikan
    # longgar. Baris kartu (SATU tinggi utk semua kartu di baris ini, dari deskripsi TERPANJANG)
    # lalu diposisikan di TENGAH zona `h`, supaya sisa ruang kosong terbagi rata atas & bawah.
    max_desc_h_in = max(
        (_estimate_wrapped_height_in(desc, 10.5, text_w_in) for _, _, _, desc in items),
        default=0,
    )
    content_h_in = 1.77 + max_desc_h_in + pad_in
    card_h_in = min(max(content_h_in, 2.2), Emu(h).inches)
    card_h = Inches(card_h_in)
    row_y = y + Inches(max(0.0, (Emu(h).inches - card_h_in) / 2))

    for idx, (num, title, stat, desc) in enumerate(items):
        cx = x + idx * (card_w + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, row_y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = t["main"]
        card.line.color.rgb = t["light"]
        card.line.width = Pt(0.75)
        _no_shadow(card)
        _modest_corner(card)
        add_badge_circle(slide, cx + pad, row_y + pad, Inches(0.42), num, t["light"], font_size=Pt(15))
        title_box = slide.shapes.add_textbox(cx + pad, row_y + pad + Inches(0.55), card_w - pad * 2, Inches(0.55))
        ttf = title_box.text_frame
        ttf.word_wrap = True
        tp = ttf.paragraphs[0]
        tp.text = title
        _set_font(tp, BODY_FONT, Pt(15), bold=True, color=WHITE)
        stat_box = slide.shapes.add_textbox(cx + pad, row_y + pad + Inches(1.1), card_w - pad * 2, Inches(0.35))
        sp = stat_box.text_frame.paragraphs[0]
        sp.text = stat
        _set_font(sp, BODY_FONT, Pt(13), bold=True, color=t["light"])
        desc_box = slide.shapes.add_textbox(cx + pad, row_y + pad + Inches(1.55), card_w - pad * 2, card_h - Inches(1.55) - pad * 2)
        dtf = desc_box.text_frame
        dtf.word_wrap = True
        dp = dtf.paragraphs[0]
        dp.text = desc
        _set_font(dp, BODY_FONT, Pt(10.5), color=RGBColor(0xE8, 0xEC, 0xE6))
    return row_y + card_h


def add_asset_ranked_bars(slide, x, y, w, items, row_h=Inches(1.0), theme: dict | None = None):
    """Alternatif visual KETIGA (selain add_asset_card_row/add_podium_row) — daftar entitas
    berperingkat dengan batang proporsional horizontal per item (badge nomor + nama + batang
    + angka) — titik variasi tampilan tambahan utk asset_cards (lihat `asset_style`). items:
    list of dict {num,name,stat,count}. Dipakai utk jumlah item BERAPA PUN (podium hanya
    cocok tepat 3)."""
    t = theme or THEME_PALETTES["green"]
    max_count = max((it.get("count") or 0) for it in items) or 1
    badge_d = Inches(0.4)
    track_h = Inches(0.16)
    stat_w = Inches(1.1)
    track_x = x + badge_d + Inches(0.18)
    track_w = w - badge_d - Inches(0.18) - stat_w - Inches(0.1)
    cur_y = y
    for it in items:
        frac = (it.get("count") or 0) / max_count if max_count else 0
        frac = max(frac, 0.04)
        add_badge_circle(slide, x, cur_y, badge_d, it["num"], t["light"], font_size=Pt(13))
        name_box = slide.shapes.add_textbox(track_x, cur_y - Inches(0.03), track_w, Inches(0.32))
        np_ = name_box.text_frame.paragraphs[0]
        np_.text = it["name"]
        _set_font(np_, BODY_FONT, Pt(13), bold=True, color=WHITE)
        stat_box = slide.shapes.add_textbox(x + w - stat_w, cur_y - Inches(0.03), stat_w, Inches(0.32))
        sp = stat_box.text_frame.paragraphs[0]
        sp.text = it["stat"]
        sp.alignment = PP_ALIGN.RIGHT
        _set_font(sp, BODY_FONT, Pt(12), bold=True, color=t["light"])

        track_y = cur_y + Inches(0.36)
        track = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, track_x, track_y, track_w, track_h)
        track.fill.solid()
        track.fill.fore_color.rgb = t["chart"]
        track.line.fill.background()
        _no_shadow(track)
        _modest_corner(track, frac=0.5)
        fill_w = Inches(max(Emu(track_w).inches * frac, 0.15))
        fill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, track_x, track_y, fill_w, track_h)
        fill.fill.solid()
        fill.fill.fore_color.rgb = t["light"]
        fill.line.fill.background()
        _no_shadow(fill)
        _modest_corner(fill, frac=0.5)
        cur_y += row_h
    return cur_y


def add_podium_row(slide, x, y, w, h, items, theme: dict | None = None):
    """items: TEPAT 3 dict {"num","name","stat"} — analog `_podium_row()` di export_pdf.py.
    Titik variasi tampilan utk asset_cards (lihat `asset_style` di generate_ppt_report),
    alternatif dari `add_asset_card_row` (baris kartu rata) — rank #1 di TENGAH & PALING
    TINGGI, meniru podium juara, dipakai kalau kebetulan item persis 3 (ranking top-3)."""
    t = theme or THEME_PALETTES["green"]
    if len(items) != 3:
        return Emu(y).inches + Emu(h).inches
    order = [1, 0, 2]  # tampil sbg [rank2, rank1, rank3] spy rank1 di tengah
    ranked = [items[i] for i in order]
    height_frac = [0.73, 1.0, 0.55]
    colors = [t["main"], t["light"], t["chart"]]
    gap = Inches(0.35)
    col_w = (w - gap * 2) / 3
    label_zone_in = 1.0
    base_bottom_in = Emu(y).inches + Emu(h).inches
    pedestal_max_h_in = max(0.8, Emu(h).inches - label_zone_in)

    for idx, item in enumerate(ranked):
        cx = x + idx * (col_w + gap)
        ped_h_in = max(0.55, pedestal_max_h_in * height_frac[idx])
        ped_top_in = base_bottom_in - ped_h_in

        name_box = slide.shapes.add_textbox(cx, Inches(ped_top_in - label_zone_in), col_w, Inches(0.32))
        np_ = name_box.text_frame.paragraphs[0]
        np_.text = item["name"]
        np_.alignment = PP_ALIGN.CENTER
        _set_font(np_, BODY_FONT, Pt(13), bold=True, color=TEXT_DARK)

        stat_box = slide.shapes.add_textbox(cx, Inches(ped_top_in - label_zone_in + 0.34), col_w, Inches(0.4))
        stf = stat_box.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.text = item["stat"]
        sp.alignment = PP_ALIGN.CENTER
        _set_font(sp, BODY_FONT, Pt(11), color=GRAY_TEXT)

        pedestal = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, Inches(ped_top_in), col_w, Inches(ped_h_in))
        pedestal.fill.solid()
        pedestal.fill.fore_color.rgb = colors[idx]
        pedestal.line.fill.background()
        _no_shadow(pedestal)
        _modest_corner(pedestal, frac=0.1)
        tf = pedestal.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pnum = tf.paragraphs[0]
        pnum.text = item["num"]
        pnum.alignment = PP_ALIGN.CENTER
        _set_font(pnum, TITLE_FONT, Pt(28), bold=True, color=WHITE)
    return base_bottom_in


def add_recommendation_timeline(slide, x, y, w, h, items, theme: dict | None = None):
    """items: 2-6 dict {"num","title","detail"} — analog `_timeline_html()` di
    export_pdf.py. Titik variasi tampilan utk recommendations (lihat `recommendation_style`
    di generate_ppt_report), alternatif dari grid kartu — garis horizontal di tengah `h`
    dengan node bergantian di ATAS/BAWAH garis, dipakai kalau jumlah item pas 2-6 (bukan
    grid biasa)."""
    t = theme or THEME_PALETTES["green"]
    n = len(items)
    if not (2 <= n <= 6):
        return False
    mid_y = y + h // 2
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, mid_y - Pt(1), w, Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = t["light"]
    line.line.fill.background()
    _no_shadow(line)

    node_d = Inches(0.4)
    slot_w = w // n
    text_w = slot_w - Inches(0.25)
    box_h = Inches(1.5)
    for idx, item in enumerate(items):
        node_cx = x + slot_w * idx + slot_w // 2
        above = (idx % 2 == 0)
        node = slide.shapes.add_shape(MSO_SHAPE.OVAL, node_cx - node_d // 2, mid_y - node_d // 2, node_d, node_d)
        node.fill.solid()
        node.fill.fore_color.rgb = t["main"]
        node.line.color.rgb = t["light"]
        node.line.width = Pt(1.5)
        _no_shadow(node)
        ntf = node.text_frame
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ntf.margin_left = ntf.margin_right = ntf.margin_top = ntf.margin_bottom = 0
        np_ = ntf.paragraphs[0]
        np_.text = item["num"]
        np_.alignment = PP_ALIGN.CENTER
        _set_font(np_, BODY_FONT, Pt(13), bold=True, color=WHITE)

        text_x = node_cx - text_w // 2
        if above:
            box_y = mid_y - node_d // 2 - box_h - Inches(0.1)
        else:
            box_y = mid_y + node_d // 2 + Inches(0.1)
        box = slide.shapes.add_textbox(text_x, box_y, text_w, box_h)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM if above else MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.text = item["title"]
        p.alignment = PP_ALIGN.CENTER
        _set_font(p, BODY_FONT, Pt(12.5), bold=True, color=TEXT_DARK)
        if item.get("detail"):
            p2 = tf.add_paragraph()
            p2.text = item["detail"]
            p2.alignment = PP_ALIGN.CENTER
            _set_font(p2, BODY_FONT, Pt(9.5), color=GRAY_TEXT)
            p2.space_before = Pt(3)
    return True


def add_recommendation_banner_list(slide, x, y, w, items, title_pt=13, detail_pt=10.5, max_y=None, theme: dict | None = None):
    """Alternatif visual KETIGA (selain grid kartu/add_recommendation_timeline) — daftar
    rekomendasi sebagai banner selebar `w` bertumpuk vertikal (badge nomor + judul + detail)
    — titik variasi tampilan tambahan utk recommendations (lihat `recommendation_style`),
    dipakai utk jumlah item BERAPA PUN (timeline dibatasi 2-6). `max_y` opsional — kalau
    total tinggi bakal melebihi, font & tinggi banner dikecilkan proporsional (pola sama
    dengan add_badge_list/add_priority_panel di file ini)."""
    t = theme or THEME_PALETTES["green"]
    if not items:
        return y
    badge_d = Inches(0.4)
    pad = Inches(0.18)
    gap = Inches(0.14)
    text_x = x + badge_d + Inches(0.35)
    text_w_in = Emu(w - badge_d - Inches(0.35) - Inches(0.15)).inches

    def _row_h_in(it, tpt, dpt):
        h = _estimate_wrapped_height_in(it["title"], tpt, text_w_in)
        if it.get("detail"):
            h += _estimate_wrapped_height_in(it["detail"], dpt, text_w_in) + 0.06
        return max(0.62, h + Emu(pad).inches * 2)

    scale = 1.0
    if max_y is not None:
        available_in = Emu(max_y - y).inches
        est_total_in = sum(_row_h_in(it, title_pt, detail_pt) + Emu(gap).inches for it in items) - Emu(gap).inches
        if available_in > 0 and est_total_in > available_in:
            scale = max(available_in / est_total_in, 0.6)

    title_pt_s = title_pt * scale
    detail_pt_s = detail_pt * scale
    cur_y = y
    for it in items:
        row_h_in = _row_h_in(it, title_pt_s, detail_pt_s)
        row_h = Inches(row_h_in)
        banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, cur_y, w, row_h)
        banner.fill.solid()
        banner.fill.fore_color.rgb = IVORY
        banner.line.color.rgb = PANEL_BORDER
        banner.line.width = Pt(0.75)
        _no_shadow(banner)
        _modest_corner(banner)
        add_badge_circle(slide, x + pad, cur_y + pad, Inches(0.4 * scale), it["num"], t["light"], font_size=Pt(13 * scale))
        # BUG YANG DIPERBAIKI: title_box dulu tinggi TETAP 0.4in, sementara det_box di bawahnya
        # diposisikan berdasarkan estimasi tinggi KONTEN sebenarnya (title_h_in, sering < 0.4in
        # utk judul 1 baris) — box tetap 0.4in penuh jadi bounding-box-nya menimpa det_box
        # (pola sama yang sudah diperbaiki di beberapa tempat lain file ini). Tinggi title_box
        # sekarang eksplisit mengikuti estimasi yang sama dipakai utk memposisikan det_box.
        title_h_in = _estimate_wrapped_height_in(it["title"], title_pt_s, text_w_in)
        title_box = slide.shapes.add_textbox(text_x, cur_y + pad - Inches(0.02), Inches(text_w_in), Inches(title_h_in + 0.08))
        ttf = title_box.text_frame
        ttf.word_wrap = True
        tp = ttf.paragraphs[0]
        tp.text = it["title"]
        _set_font(tp, BODY_FONT, Pt(title_pt_s), bold=True, color=TEXT_DARK)
        if it.get("detail"):
            det_box = slide.shapes.add_textbox(text_x, cur_y + pad + Inches(title_h_in + 0.04), Inches(text_w_in), Inches(max(0.3, row_h_in - Emu(pad).inches - title_h_in)))
            dtf = det_box.text_frame
            dtf.word_wrap = True
            dp = dtf.paragraphs[0]
            dp.text = it["detail"]
            _set_font(dp, BODY_FONT, Pt(detail_pt_s), color=GRAY_TEXT)
        cur_y += row_h + gap
    return cur_y


def add_native_bar_chart(slide, x, y, cx, cy, categories, values, colors=None, horizontal=False):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Jumlah", values)
    chart_type = XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED
    gframe = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data)
    chart = gframe.chart
    chart.has_legend = False
    # BUG YANG DIPERBAIKI (dilaporkan user, disertai tangkapan layar): python-pptx/PowerPoint
    # menampilkan judul chart DEFAULT berupa nama series ("Jumlah") kalau tidak dimatikan
    # eksplisit — nongol polos warna hitam, tidak nyambung sama gaya desain sekitarnya.
    chart.has_title = False

    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.size = Pt(11)
    dl.font.name = BODY_FONT
    dl.font.bold = True
    try:
        dl.number_format = "0"
        dl.number_format_is_linked = False
    except Exception:
        pass

    series = plot.series[0]
    if colors:
        for i, pt in enumerate(series.points):
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = colors[i % len(colors)]
    else:
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = GREEN_MAIN

    try:
        chart.category_axis.has_major_gridlines = False
        chart.value_axis.has_major_gridlines = False
        chart.category_axis.tick_labels.font.size = Pt(11)
        chart.category_axis.tick_labels.font.name = BODY_FONT
        chart.value_axis.visible = False
    except Exception:
        pass
    return gframe


def add_native_doughnut_chart(slide, x, y, cx, cy, categories, values, colors=None, theme: dict | None = None):
    """Alternatif visual utk distribusi kategori (selain add_native_bar_chart) — titik
    variasi tampilan antar generate (lihat `category_style` di generate_ppt_report), chart
    doughnut NATIVE PowerPoint (bukan gambar statis) supaya tetap bisa diedit user kalau mau,
    konsisten dengan seluruh chart lain di file ini."""
    t = theme or THEME_PALETTES["green"]
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Jumlah", values)
    gframe = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data)
    chart = gframe.chart
    chart.has_legend = False
    chart.has_title = False  # lihat catatan sama di add_native_bar_chart soal judul default "Jumlah"

    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.size = Pt(11)
    dl.font.name = BODY_FONT
    dl.font.bold = True
    dl.font.color.rgb = WHITE
    try:
        dl.number_format = "0"
        dl.number_format_is_linked = False
    except Exception:
        pass

    series = plot.series[0]
    palette = colors or CATEGORY_COLOR_RAMP
    for i, pt in enumerate(series.points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = palette[i % len(palette)]
        pt.format.line.color.rgb = t["bg"]
        pt.format.line.width = Pt(1.5)
    return gframe


def add_stacked_proportion_bar(slide, x, y, w, values, colors=None, height=Inches(0.5)):
    """Alternatif visual KETIGA (selain add_native_bar_chart/add_native_doughnut_chart) — satu
    batang penuh dibagi proporsional per kategori (gaya "100% stacked bar") — titik variasi
    tampilan tambahan supaya laporan tidak melulu bar-per-baris atau donut (lihat
    category_style/status_style di generate_ppt_report). DIPASANGKAN dengan panel legend
    TERPISAH DI BAWAHNYA oleh pemanggil (bukan di samping) — batang cuma setinggi `height`,
    kalau dipasangkan sejajar dengan panel legend yang jauh lebih tinggi bakal menyisakan
    ruang kosong besar, persis kelas masalah "ruang kosong tidak proporsional" yang sudah
    diperbaiki di tempat lain di file ini."""
    total = sum(values) or 1
    w_in = Emu(w).inches
    cur_x_in = Emu(x).inches
    for i, val in enumerate(values):
        frac = (val / total) if total else 0
        seg_w_in = w_in * frac
        if seg_w_in <= 0:
            continue
        color = colors[i] if colors else CATEGORY_COLOR_RAMP[i % len(CATEGORY_COLOR_RAMP)]
        seg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cur_x_in), y, Inches(seg_w_in), height)
        seg.fill.solid()
        seg.fill.fore_color.rgb = color
        seg.line.color.rgb = WHITE
        seg.line.width = Pt(1)
        _no_shadow(seg)
        cur_x_in += seg_w_in
    return Emu(y).inches + Emu(height).inches


def add_native_table(slide, x, y, w, h, headers, rows, highlight_indices=None, theme: dict | None = None):
    t = theme or THEME_PALETTES["green"]
    highlight_indices = highlight_indices or set()
    n_rows = len(rows) + 1
    n_cols = len(headers)
    gframe = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    table = gframe.table

    for c, htext in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = t["bg"]
        cell.margin_left = cell.margin_right = Inches(0.08)
        p = cell.text_frame.paragraphs[0]
        p.text = htext
        _set_font(p, BODY_FONT, Pt(11), bold=True, color=WHITE)

    for r, row_vals in enumerate(rows):
        is_open = r in highlight_indices
        for c, val in enumerate(row_vals):
            cell = table.cell(r + 1, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RED_CRIT_BG if is_open else (IVORY if r % 2 == 0 else WHITE)
            cell.margin_left = cell.margin_right = Inches(0.08)
            p = cell.text_frame.paragraphs[0]
            p.text = str(val)
            is_status_col = c == n_cols - 1
            _set_font(
                p, BODY_FONT, Pt(10.5),
                bold=bool(is_open and is_status_col),
                color=(RED_CRIT if (is_open and is_status_col) else TEXT_DARK),
            )
    return gframe


def add_split_cover_slide(prs, block, flourish_corner, logo_path, theme: dict | None = None):
    """Varian cover 2-kolom warna penuh (emas kiri + hijau kanan, angka hero besar di kolom
    emas) — analog `_split_cover_td()` di export_pdf.py, titik variasi tampilan (lihat
    `cover_style` di generate_ppt_report), alternatif dari cover 1-warna standar."""
    t = theme or THEME_PALETTES["green"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left_w = Inches(13.33 * 0.37)
    right_x = left_w
    right_w = SLIDE_W - left_w
    _fill_rect_bg(slide, 0, 0, left_w, SLIDE_H, t["light"])
    _fill_rect_bg(slide, right_x, 0, right_w, SLIDE_H, t["bg"])
    add_corner_flourish(slide, flourish_corner, area_x=right_x, area_w=right_w, theme=t)
    add_logo(slide, logo_path)

    value, label = block.get("hero_stat") or (str(block.get("total_records", "")), "Total Data")
    hero_kicker = block.get("hero_stat_kicker", "CAPAIAN")

    kicker_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.5), left_w - Inches(0.7), Inches(0.3))
    kp = kicker_box.text_frame.paragraphs[0]
    kp.text = hero_kicker.upper()
    _set_font(kp, BODY_FONT, Pt(10), bold=True, color=TEXT_DARK)

    value_w_in = Emu(left_w - Inches(0.6)).inches
    value_h_in = _estimate_wrapped_height_in(str(value), 54, value_w_in)
    value_box = slide.shapes.add_textbox(Inches(0.4), Inches(3.15), left_w - Inches(0.6), Inches(value_h_in + 0.1))
    vtf = value_box.text_frame
    vtf.word_wrap = True
    vp = vtf.paragraphs[0]
    vp.text = str(value)
    _set_font(vp, TITLE_FONT, Pt(54), bold=True, color=t["bg"])

    label_top_in = 3.15 + value_h_in + 0.15
    label_box = slide.shapes.add_textbox(Inches(0.45), Inches(label_top_in), left_w - Inches(0.7), Inches(0.4))
    lp = label_box.text_frame.paragraphs[0]
    lp.text = label
    _set_font(lp, BODY_FONT, Pt(12), color=TEXT_DARK)

    footer_box = slide.shapes.add_textbox(Inches(0.45), SLIDE_H - Inches(0.7), left_w - Inches(0.7), Inches(0.3))
    fp = footer_box.text_frame.paragraphs[0]
    fp.text = block["header_title"]
    _set_font(fp, BODY_FONT, Pt(10), bold=True, color=TEXT_DARK)

    title_text = block["title"]
    if len(title_text) > 55:
        title_size_pt = 26
    elif len(title_text) > 40:
        title_size_pt = 30
    elif len(title_text) > 28:
        title_size_pt = 36
    else:
        title_size_pt = 42

    text_x = right_x + Inches(0.5)
    text_w = right_w - Inches(1.0)

    kicker2_box = slide.shapes.add_textbox(text_x, Inches(2.15), text_w, Inches(0.3))
    kp2 = kicker2_box.text_frame.paragraphs[0]
    kp2.text = block["kicker"].upper()
    _set_font(kp2, BODY_FONT, Pt(10.5), bold=True, color=t["light"])

    title_top_in = 2.5
    title_height_in = _estimate_wrapped_height_in(title_text, title_size_pt, Emu(text_w).inches)
    title_box = slide.shapes.add_textbox(text_x, Inches(title_top_in), text_w, Inches(title_height_in + 0.1))
    ttf = title_box.text_frame
    ttf.word_wrap = True
    tp = ttf.paragraphs[0]
    tp.text = title_text
    _set_font(tp, TITLE_FONT, Pt(title_size_pt), bold=True, color=WHITE)

    sub_top_in = max(title_top_in + title_height_in + 0.15, 3.4)
    sub_box = slide.shapes.add_textbox(text_x, Inches(sub_top_in), text_w, Inches(0.5))
    sp = sub_box.text_frame.paragraphs[0]
    sp.text = block["subtitle"]
    _set_font(sp, BODY_FONT, Pt(14), color=WHITE)

    info_top_in = max(sub_top_in + 0.6, 4.2)
    info_box = slide.shapes.add_textbox(text_x, Inches(info_top_in), text_w, Inches(0.9))
    itf = info_box.text_frame
    itf.word_wrap = True
    p1 = itf.paragraphs[0]
    p1.text = f'{block["period_label"]} {block["period_text"]}'
    _set_font(p1, BODY_FONT, Pt(12), color=WHITE)
    p2 = itf.add_paragraph()
    p2.text = block["info_line"]
    _set_font(p2, BODY_FONT, Pt(12), color=t["soft"])
    p2.space_before = Pt(6)

    return slide


def add_split_closing_slide(prs, block, flourish_corner, theme: dict | None = None):
    """Varian penutup berpasangan dgn `add_split_cover_slide` — angka hero yang SAMA
    ditampilkan lagi di kolom emas kiri (mengulang temuan utama, gaya "bookend" laporan
    eksekutif), analog `_split_closing_td()` di export_pdf.py."""
    t = theme or THEME_PALETTES["green"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left_w = Inches(13.33 * 0.37)
    right_x = left_w
    right_w = SLIDE_W - left_w
    _fill_rect_bg(slide, 0, 0, left_w, SLIDE_H, t["light"])
    _fill_rect_bg(slide, right_x, 0, right_w, SLIDE_H, t["bg"])
    add_corner_flourish(slide, flourish_corner, area_x=right_x, area_w=right_w, theme=t)

    value, label = block.get("hero_stat") or ("", "")
    value_w_in = Emu(left_w - Inches(0.6)).inches
    value_h_in = _estimate_wrapped_height_in(str(value), 42, value_w_in)
    value_box = slide.shapes.add_textbox(Inches(0.4), Inches(3.15), left_w - Inches(0.6), Inches(value_h_in + 0.1))
    vtf = value_box.text_frame
    vtf.word_wrap = True
    vp = vtf.paragraphs[0]
    vp.text = str(value)
    _set_font(vp, TITLE_FONT, Pt(42), bold=True, color=t["bg"])

    label_top_in = 3.15 + value_h_in + 0.12
    label_box = slide.shapes.add_textbox(Inches(0.45), Inches(label_top_in), left_w - Inches(0.7), Inches(0.4))
    lp = label_box.text_frame.paragraphs[0]
    lp.text = label
    _set_font(lp, BODY_FONT, Pt(11), color=TEXT_DARK)

    text_x = right_x + Inches(0.5)
    text_w = right_w - Inches(1.0)
    thank_you_h_in = _estimate_wrapped_height_in(block["thank_you"], 38, Emu(text_w).inches)
    title_box = slide.shapes.add_textbox(text_x, Inches(3.0), text_w, Inches(thank_you_h_in + 0.1))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = block["thank_you"]
    _set_font(tp, TITLE_FONT, Pt(38), bold=True, color=WHITE)

    sub_top_in = max(3.0 + thank_you_h_in + 0.15, 3.85)
    sub_box = slide.shapes.add_textbox(text_x, Inches(sub_top_in), text_w, Inches(0.5))
    stf = sub_box.text_frame
    stf.word_wrap = True
    sp = stf.paragraphs[0]
    sp.text = block["title"]
    _set_font(sp, BODY_FONT, Pt(13), color=WHITE)

    sub_height_in = _estimate_wrapped_height_in(block["title"], 13, Emu(text_w).inches)
    note_top_in = max(sub_top_in + sub_height_in + 0.1, 4.4)
    note_box = slide.shapes.add_textbox(text_x, Inches(note_top_in), text_w, Inches(0.4))
    np_ = note_box.text_frame.paragraphs[0]
    np_.text = block["note"]
    _set_font(np_, BODY_FONT, Pt(11), italic=True, color=t["soft"])

    return slide


# ============================================================================
# Konten turunan dari data (bukan dari AI) — deterministik
# ============================================================================


@dataclass
class _PptBlockContext:
    """Kumpulan variabel variasi tampilan (dipilih SEKALI per generate, lihat
    generate_ppt_report) + report/prs/logo_path yang dibutuhkan LEBIH DARI SATU builder
    slide di bawah — dioper ke tiap builder supaya signature-nya seragam (block, ctx)
    alih-alih daftar parameter berbeda-beda per jenis slide. Sebelumnya semua builder ini
    adalah cabang if/elif di dalam SATU fungsi generate_ppt_report sepanjang ~650 baris —
    dipecah jadi fungsi terpisah (murni supaya lebih mudah dibaca/diubah 1 jenis slide
    tanpa perlu scroll baca semuanya), TIDAK ada perubahan HASIL AKHIR sama sekali
    (diverifikasi PPTX yang dihasilkan tetap konsisten sebelum & sesudah pemecahan ini).

    cover_hero_stat DIISI oleh _build_cover_slide, DIBACA oleh _build_closing_slide
    (bookend angka hero yang sama di cover & penutup saat cover_style="split") — satu-
    satunya state yang mengalir ANTAR pemanggilan builder, makanya ctx harus objek yang
    sama dioper ke semua builder dalam 1 kali generate, bukan dibuat ulang tiap slide.

    Beda dari _PdfBlockContext (export_pdf.py): builder di sini TIDAK mengembalikan HTML,
    melainkan langsung menggambar shape ke slide baru di ctx.prs (efek samping), lalu
    me-return objek Slide yang harus di-footer-nomori (None kalau slide itu cover/penutup,
    yang memang selalu dikecualikan dari penomoran — lihat _PPT_BLOCK_BUILDERS di bawah).
    """
    report: Report
    prs: Presentation
    logo_path: object
    panel_side: str
    stat_cols: int
    card_cols: int
    flourish_corner: str
    accent_bar_color: RGBColor
    category_style: str
    status_style: str
    cover_style: str
    asset_style: str
    recommendation_style: str
    kicker_ringkasan: str
    kicker_analisis: str
    # Palet warna tema (report.theme_color) — 5 peran, sama persis dgn export_pdf.py. Dipakai
    # di elemen BRAND/struktural — TIDAK PERNAH di SEVERITY_COLOR/kondisional is_critical.
    accent_main: RGBColor
    accent_bg: RGBColor
    accent_chart: RGBColor
    accent_light: RGBColor
    accent_soft: RGBColor
    theme: dict | None = None
    cover_hero_stat: dict | None = None


def _build_cover_slide(block: dict, ctx: _PptBlockContext):
    ctx.cover_hero_stat = block.get("hero_stat")
    if ctx.cover_style == "split":
        add_split_cover_slide(ctx.prs, block, ctx.flourish_corner, ctx.logo_path, theme=ctx.theme)
        return None
    cover = ctx.prs.slides.add_slide(ctx.prs.slide_layouts[6])
    add_dark_bg(cover, theme=ctx.theme)
    add_corner_flourish(cover, ctx.flourish_corner, theme=ctx.theme)
    add_logo(cover, ctx.logo_path)

    add_kicker(cover, block["kicker"], color=ctx.accent_light, y=Inches(1.7))

    # Ukuran font judul menyesuaikan panjangnya (laporan bisa dari domain apa saja
    # - SOC, keuangan, KPI, dll - judulnya bisa jauh lebih panjang/pendek dari
    # contoh mana pun), lalu posisi subtitle & info DIHITUNG dari perkiraan tinggi
    # judul yang sebenarnya — bukan angka tetap — supaya tidak pernah tertimpa
    # berapa pun baris yang dibutuhkan judul untuk wrap.
    title_text = block["title"]
    if len(title_text) > 55:
        title_size_pt = 26
    elif len(title_text) > 40:
        title_size_pt = 32
    elif len(title_text) > 28:
        title_size_pt = 38
    else:
        title_size_pt = 44

    title_top_in = 2.1
    title_box = cover.shapes.add_textbox(MARGIN_X, Inches(title_top_in), Inches(9.5), Inches(2.4))
    ttf = title_box.text_frame
    ttf.word_wrap = True
    tp = ttf.paragraphs[0]
    tp.text = title_text
    _set_font(tp, TITLE_FONT, Pt(title_size_pt), bold=True, color=WHITE)

    title_height_in = _estimate_wrapped_height_in(title_text, title_size_pt, 9.5)
    sub_top_in = max(title_top_in + title_height_in + 0.15, 3.05)
    sub_box = cover.shapes.add_textbox(MARGIN_X, Inches(sub_top_in), Inches(9.5), Inches(0.5))
    sp = sub_box.text_frame.paragraphs[0]
    sp.text = block["subtitle"]
    _set_font(sp, BODY_FONT, Pt(15), color=WHITE)

    info_top_in = max(sub_top_in + 0.65, 3.9)
    info_box = cover.shapes.add_textbox(MARGIN_X, Inches(info_top_in), Inches(9.5), Inches(0.9))
    itf = info_box.text_frame
    itf.word_wrap = True
    p1 = itf.paragraphs[0]
    p1.text = f'{block["period_label"]} {block["period_text"]}'
    _set_font(p1, BODY_FONT, Pt(12.5), color=WHITE)
    p2 = itf.add_paragraph()
    p2.text = block["info_line"]
    _set_font(p2, BODY_FONT, Pt(12.5), color=ctx.accent_soft)
    p2.space_before = Pt(6)

    footer_l = cover.shapes.add_textbox(MARGIN_X, SLIDE_H - Inches(0.55), Inches(5), Inches(0.3))
    flp = footer_l.text_frame.paragraphs[0]
    flp.text = block["header_title"]
    _set_font(flp, BODY_FONT, Pt(10), bold=True, color=WHITE)
    return None


def _build_intro_slide(block: dict, ctx: _PptBlockContext):
    bg_slide = ctx.prs.slides.add_slide(ctx.prs.slide_layouts[6])
    add_logo(bg_slide, ctx.logo_path)
    add_kicker(bg_slide, block["kicker"], color=ctx.accent_main)
    title_bottom = add_title(bg_slide, block["title"])

    left_w = Inches(7.0) if ctx.panel_side == "right" else Inches(4.9)
    left_x = MARGIN_X if ctx.panel_side == "right" else MARGIN_X + Inches(4.9) + Inches(0.4)
    panel_x = MARGIN_X + Inches(7.0) + Inches(0.4) if ctx.panel_side == "right" else MARGIN_X
    panel_w = Inches(4.9)

    # content_top turun kalau judul wrap ke 2 baris (lihat catatan di add_title);
    # jarak RELATIF antar elemen di bawahnya (badge list, dst) tetap sama.
    content_top = max(title_bottom, 1.65)
    para_box = bg_slide.shapes.add_textbox(left_x, Inches(content_top), left_w, Inches(1.3))
    ptf = para_box.text_frame
    ptf.word_wrap = True
    pp = ptf.paragraphs[0]
    pp.text = block["purpose_text"]
    _set_font(pp, BODY_FONT, Pt(13), color=GRAY_TEXT)

    objectives = [(o["num"], o["title"], o["detail"]) for o in block["objectives"]]
    add_badge_list(bg_slide, left_x, Inches(content_top + 1.4), left_w, objectives, badge_color=ctx.accent_main, row_h=Inches(1.05), max_y=SLIDE_H - Inches(0.4))

    scope = block["scope"]
    scope_rows = [
        (scope["period_label"], scope["period_text"]),
        (scope["total_event_label"], scope["total_records_text"]),
        (scope["source_file_label"], scope["input_file_name"]),
        (scope["data_type_label_label"], scope["data_type_label"]),
    ]
    add_ivory_panel(
        bg_slide, panel_x, Inches(content_top), panel_w, Inches(4.3),
        "i", scope["panel_title"], scope_rows, mode="kv",
        footnote=scope["footnote"], theme=ctx.theme,
    )
    return bg_slide


def _build_executive_summary_slide(block: dict, ctx: _PptBlockContext):
    exec_slide = ctx.prs.slides.add_slide(ctx.prs.slide_layouts[6])
    add_dark_bg(exec_slide, theme=ctx.theme)
    add_logo(exec_slide, ctx.logo_path)
    add_kicker(exec_slide, ctx.kicker_ringkasan, color=ctx.accent_light)
    title_bottom = add_title(exec_slide, block["heading"], color=WHITE)

    # BUG NYATA YANG DIPERBAIKI (dilaporkan user, disertai tangkapan layar): grid
    # kartu KPI SEBELUMNYA selalu mulai di Y=1.7in tetap, menabrak baris ke-2
    # heading kalau heading-nya panjang & wrap (umum utk data non-SOC). grid_y
    # sekarang mengikuti tinggi heading sungguhan; grid_h dikurangi secukupnya
    # supaya grid tidak meluber ke luar slide kalau grid_y ikut turun.
    grid_y_in = max(title_bottom + 0.2, 1.7)
    grid_h_in = max(2.2, 4.3 - (grid_y_in - 1.7))

    # Keluhan nyata dari pengguna: kalau item KPI cuma sedikit (mis. 2 kartu, 1
    # baris), grid+caption cuma menempati sepertiga atas slide, sisanya kosong total
    # sampai footer. Perkirakan DULU tinggi grid+caption yang akan dihasilkan (tanpa
    # menggambar apapun — replikasi kecil rumus di add_stat_card_grid), lalu kalau
    # ternyata jauh lebih pendek dari ruang yang tersedia, geser blok ini ke bawah
    # supaya "kosongnya" terbagi rata di atas & bawah, bukan menumpuk di bawah saja.
    n_stat_items = len(block["stat_items"])
    if n_stat_items:
        rows_n = math.ceil(n_stat_items / ctx.stat_cols)
        max_card_h_in = {1: 2.6, 2: 2.0}.get(rows_n, 1.6)
        natural_card_h_in = (grid_h_in - 0.2 * (rows_n - 1)) / rows_n
        card_h_in = min(natural_card_h_in, max_card_h_in)
        grid_total_h_in = rows_n * card_h_in + (rows_n - 1) * 0.2
        caption_h_in = _estimate_wrapped_height_in(block["caption"], 11.5, Emu(CONTENT_W).inches)
        content_total_in = grid_total_h_in + 0.25 + caption_h_in
        available_in = 6.95 - grid_y_in
        if available_in > content_total_in:
            grid_y_in += (available_in - content_total_in) / 2

    grid_bottom = add_stat_card_grid(exec_slide, MARGIN_X, Inches(grid_y_in), CONTENT_W, Inches(grid_h_in), block["stat_items"], cols=ctx.stat_cols, dark=True, theme=ctx.theme)

    # BUG YANG DIPERBAIKI: box caption dulu tinggi TETAP 0.9in apa pun panjang
    # teksnya — kalau blok grid+caption digeser turun (lihat centering di atas),
    # box tetap 0.9in penuh bisa meluber ke zona footer walau isi captionnya cuma
    # 1 baris pendek. Tinggi box sekarang mengikuti estimasi wrap sungguhan (dibatasi
    # minimum 0.4in), sama seperti perhitungan yang dipakai utk centering di atas.
    caption_box_h_in = max(0.4, _estimate_wrapped_height_in(block["caption"], 11.5, Emu(CONTENT_W).inches) + 0.1)
    cap_box = exec_slide.shapes.add_textbox(MARGIN_X, grid_bottom + Inches(0.25), CONTENT_W, Inches(caption_box_h_in))
    ctf = cap_box.text_frame
    ctf.word_wrap = True
    cp = ctf.paragraphs[0]
    cp.text = block["caption"]
    _set_font(cp, BODY_FONT, Pt(11.5), italic=True, color=ctx.accent_soft)
    return exec_slide


def _build_dynamic_section_slide(block: dict, ctx: _PptBlockContext):
    dyn_slide = ctx.prs.slides.add_slide(ctx.prs.slide_layouts[6])
    add_logo(dyn_slide, ctx.logo_path)
    add_kicker(dyn_slide, block["kicker"], color=ctx.accent_main)
    title_bottom = add_title(dyn_slide, block["title"])
    content_top = max(title_bottom + 0.2, 1.6)

    # Panel angka/daftar di samping teks (kalau tersedia) — supaya slide narasi
    # tidak cuma "judul + 1 paragraf" mubazir ruang kosong (temuan user), dan
    # berselang-seling 2 pola (angka besar vs daftar ringkas) via layout_variant
    # yang sudah ditentukan report_render_logic.py.
    has_aux = bool(block.get("aux_stat") or block.get("aux_list"))
    text_w = Inches(7.3) if has_aux else CONTENT_W
    # panel_side menentukan teks di kiri+panel di kanan (default) atau dibalik —
    # titik variasi tampilan yang sama dipakai halaman lain (lihat komentar atas).
    if has_aux and ctx.panel_side == "left":
        panel_w = CONTENT_W - text_w - Inches(0.4)
        panel_x = MARGIN_X
        text_x = panel_x + panel_w + Inches(0.4)
    else:
        text_x = MARGIN_X
        panel_x = MARGIN_X + text_w + Inches(0.4)
        panel_w = SLIDE_W - MARGIN_X - panel_x
    text_box = dyn_slide.shapes.add_textbox(text_x, Inches(content_top), text_w, Inches(4.9))
    ttf3 = text_box.text_frame
    ttf3.word_wrap = True
    tp4 = ttf3.paragraphs[0]
    tp4.text = block["text"]
    _set_font(tp4, BODY_FONT, Pt(13), color=GRAY_TEXT)

    if has_aux:
        if block.get("aux_stat"):
            value, label = block["aux_stat"]
            add_critical_highlight_panel(dyn_slide, panel_x, Inches(content_top), panel_w, Inches(2.6), value, label, theme=ctx.theme)
        else:
            rows = [(it["label"], it["value"]) for it in block["aux_list"]]
            panel_title = "Data Highlight" if is_english(ctx.report) else "Sorotan Data"
            add_ivory_panel(dyn_slide, panel_x, Inches(content_top), panel_w, Inches(3.4), "i", panel_title, rows, mode="kv", theme=ctx.theme)
    return dyn_slide


def _build_category_distribution_slide(block: dict, ctx: _PptBlockContext):
    cat_slide = ctx.prs.slides.add_slide(ctx.prs.slide_layouts[6])
    add_logo(cat_slide, ctx.logo_path)
    add_kicker(cat_slide, ctx.kicker_analisis, color=ctx.accent_main)
    title_bottom = add_title(cat_slide, block["title"])

    intro_y = max(title_bottom + 0.15, 1.45)
    intro_box = cat_slide.shapes.add_textbox(MARGIN_X, Inches(intro_y), CONTENT_W, Inches(0.5))
    itf2 = intro_box.text_frame
    itf2.word_wrap = True
    ip = itf2.paragraphs[0]
    ip.text = block["intro"]
    _set_font(ip, BODY_FONT, Pt(12), color=GRAY_TEXT)

    body_y = intro_y + 0.65
    cat_has_caption = bool(block.get("ai_caption"))
    cat_body_h = Inches(4.0) if cat_has_caption else Inches(4.6)
    # Ramp warna kategori/status DITURUNKAN dari tema (report.theme_color), bukan konstanta
    # hijau/emas tetap — sama seperti export_pdf.py. GRAY_TEXT tetap warna ke-5 (netral).
    ramp = [ctx.accent_main, ctx.accent_chart, ctx.accent_light, ctx.accent_soft, GRAY_TEXT]
    legend_rows = [
        (ramp[l["color_index"] % len(ramp)], l["name"], f'{l["pct"]}%')
        for l in block["legend"]
    ]
    # Titik variasi tampilan: bar horizontal (warna accent hijau/emas gantian),
    # donut chart NATIVE PowerPoint, ATAU batang proporsi 100% bersegmen (lihat
    # category_style/accent_bar_color di atas) — datanya identik, cuma cara
    # visualnya beda tiap generate.
    cat_content_bottom_in = body_y + Emu(cat_body_h).inches
    if ctx.category_style == "stacked":
        # Ditumpuk VERTIKAL (batang pendek lalu panel legend penuh di bawahnya),
        # BUKAN dipasangkan sejajar seperti bar/donut — batang cuma ~0.55in, kalau
        # dipaksa sejajar dgn panel setinggi cat_body_h bakal menyisakan ruang
        # kosong besar di sampingnya (kelas masalah yang sudah diperbaiki di
        # tempat lain di file ini). Tinggi panel legend di sini SEKARANG mengikuti
        # jumlah baris sebenarnya (bukan sisa ruang cat_body_h yang longgar) — panel
        # add_ivory_panel PPT tidak menyusut sendiri ke kontennya (beda dari versi
        # PDF), jadi kalau dikasih tinggi generus, hasilnya kartu besar nyaris kosong.
        bar_h = Inches(0.55)
        seg_colors = [ramp[l["color_index"] % len(ramp)] for l in block["legend"]]
        add_stacked_proportion_bar(cat_slide, MARGIN_X, Inches(body_y), CONTENT_W, block["values"], colors=seg_colors, height=bar_h)
        legend_y_in = body_y + Emu(bar_h).inches + 0.35
        legend_h_in = 0.85 + 0.4 * len(legend_rows) + (0.35 if block["footnote"] else 0)
        add_ivory_panel(
            cat_slide, MARGIN_X, Inches(legend_y_in), CONTENT_W, Inches(legend_h_in),
            "%", block["legend_panel_title"], legend_rows, mode="legend",
            footnote=block["footnote"], theme=ctx.theme,
        )
        cat_content_bottom_in = legend_y_in + legend_h_in
    else:
        chart_w = Inches(7.3) if ctx.panel_side == "right" else Inches(4.9)
        chart_x = MARGIN_X if ctx.panel_side == "right" else MARGIN_X + Inches(4.9) + Inches(0.4)
        panel_x2 = MARGIN_X + Inches(7.3) + Inches(0.4) if ctx.panel_side == "right" else MARGIN_X
        if ctx.category_style == "donut":
            donut_side = min(Emu(chart_w).inches, Emu(cat_body_h).inches)
            donut_x = chart_x + Inches((Emu(chart_w).inches - donut_side) / 2)
            donut_y = Inches(body_y + (Emu(cat_body_h).inches - donut_side) / 2)
            add_native_doughnut_chart(
                cat_slide, donut_x, donut_y, Inches(donut_side), Inches(donut_side),
                block["categories"], block["values"],
                colors=[ramp[l["color_index"] % len(ramp)] for l in block["legend"]], theme=ctx.theme,
            )
        else:
            add_native_bar_chart(
                cat_slide, chart_x, Inches(body_y), chart_w, cat_body_h,
                list(reversed(block["categories"])), list(reversed(block["values"])),
                horizontal=True, colors=[ctx.accent_bar_color],
            )
        add_ivory_panel(
            cat_slide, panel_x2, Inches(body_y), Inches(4.9), cat_body_h,
            "%", block["legend_panel_title"], legend_rows, mode="legend",
            footnote=block["footnote"], theme=ctx.theme,
        )
    if cat_has_caption:
        add_ai_insight_strip(cat_slide, MARGIN_X, Inches(cat_content_bottom_in + 0.12), CONTENT_W, block["ai_caption"])
    return cat_slide


def _build_severity_distribution_slide(block: dict, ctx: _PptBlockContext):
    sev_slide = ctx.prs.slides.add_slide(ctx.prs.slide_layouts[6])
    add_logo(sev_slide, ctx.logo_path)
    add_kicker(sev_slide, ctx.kicker_analisis, color=ctx.accent_main)
    title_bottom = add_title(sev_slide, block["title"])

    intro_y = max(title_bottom + 0.15, 1.45)
    intro_box = sev_slide.shapes.add_textbox(MARGIN_X, Inches(intro_y), CONTENT_W, Inches(0.5))
    itf3 = intro_box.text_frame
    itf3.word_wrap = True
    ip3 = itf3.paragraphs[0]
    ip3.text = block["intro"]
    _set_font(ip3, BODY_FONT, Pt(12), color=GRAY_TEXT)

    sev_body_y = intro_y + 0.65
    sev_has_caption = bool(block.get("ai_caption"))
    sev_body_h = Inches(4.0) if sev_has_caption else Inches(4.6)
    sev_colors = [SEVERITY_COLOR[k] for k in block["severity_keys"]]
    # panel_side: chart di kiri+panel di kanan (default) atau dibalik — ukuran
    # kolom tetap sama, cuma posisinya ditukar.
    if ctx.panel_side == "left":
        sev_panel_x, sev_chart_x = MARGIN_X, MARGIN_X + Inches(4.7)
    else:
        sev_chart_x, sev_panel_x = MARGIN_X, MARGIN_X + Inches(8.5)
    add_native_bar_chart(sev_slide, sev_chart_x, Inches(sev_body_y), Inches(8.1), sev_body_h, block["categories"], block["values"], colors=sev_colors)

    add_critical_highlight_panel(
        sev_slide, sev_panel_x, Inches(sev_body_y), Inches(4.3), sev_body_h,
        f'{block["crit_pct"]}%', block["panel_text"], block["detail_text"], theme=ctx.theme,
    )
    if sev_has_caption:
        add_ai_insight_strip(sev_slide, MARGIN_X, Inches(sev_body_y) + sev_body_h + Inches(0.12), CONTENT_W, block["ai_caption"])
    return sev_slide


def _build_status_distribution_slide(block: dict, ctx: _PptBlockContext):
    status_slide = ctx.prs.slides.add_slide(ctx.prs.slide_layouts[6])
    add_logo(status_slide, ctx.logo_path)
    add_kicker(status_slide, ctx.kicker_analisis, color=ctx.accent_main)
    title_bottom = add_title(status_slide, block["title"])

    intro_y = max(title_bottom + 0.15, 1.45)
    intro_box = status_slide.shapes.add_textbox(MARGIN_X, Inches(intro_y), CONTENT_W, Inches(0.5))
    itf4 = intro_box.text_frame
    itf4.word_wrap = True
    ip4 = itf4.paragraphs[0]
    ip4.text = block["intro"]
    _set_font(ip4, BODY_FONT, Pt(12), color=GRAY_TEXT)

    status_body_y = intro_y + 0.65
    status_has_caption = bool(block.get("ai_caption"))
    status_body_h = Inches(4.0) if status_has_caption else Inches(4.6)
    status_content_bottom_in = status_body_y + Emu(status_body_h).inches
    # Titik variasi tampilan (independen dari category_style — lihat status_style di
    # atas): donut/stacked butuh panel legend (segmen tidak ber-label nama sendiri,
    # beda dari bar chart yang sumbu kategorinya otomatis jadi label) — dipola sama
    # persis seperti category_distribution.
    ramp = [ctx.accent_main, ctx.accent_chart, ctx.accent_light, ctx.accent_soft, GRAY_TEXT]
    if ctx.status_style in ("donut", "stacked"):
        status_total = sum(block["values"]) or 1
        status_colors = [ramp[i % len(ramp)] for i in range(len(block["values"]))]
        status_legend_rows = [
            (status_colors[i], name, f"{round(val / status_total * 100, 1)}%")
            for i, (name, val) in enumerate(zip(block["categories"], block["values"]))
        ]
        status_legend_title = "Status Proportion" if is_english(ctx.report) else "Proporsi Status"
        if ctx.status_style == "donut":
            status_chart_w = Inches(7.3) if ctx.panel_side == "right" else Inches(4.9)
            status_chart_x = MARGIN_X if ctx.panel_side == "right" else MARGIN_X + Inches(4.9) + Inches(0.4)
            status_panel_x = MARGIN_X + Inches(7.3) + Inches(0.4) if ctx.panel_side == "right" else MARGIN_X
            donut_side = min(Emu(status_chart_w).inches, Emu(status_body_h).inches)
            donut_x = status_chart_x + Inches((Emu(status_chart_w).inches - donut_side) / 2)
            donut_y = Inches(status_body_y + (Emu(status_body_h).inches - donut_side) / 2)
            add_native_doughnut_chart(
                status_slide, donut_x, donut_y, Inches(donut_side), Inches(donut_side),
                block["categories"], block["values"], colors=status_colors, theme=ctx.theme,
            )
            add_ivory_panel(
                status_slide, status_panel_x, Inches(status_body_y), Inches(4.9), status_body_h,
                "%", status_legend_title, status_legend_rows, mode="legend", theme=ctx.theme,
            )
        else:
            # "stacked" ditumpuk vertikal — lihat catatan sama di category_distribution
            # soal kenapa tinggi panel mengikuti jumlah baris, bukan sisa ruang.
            bar_h = Inches(0.55)
            add_stacked_proportion_bar(status_slide, MARGIN_X, Inches(status_body_y), CONTENT_W, block["values"], colors=status_colors, height=bar_h)
            legend_y_in = status_body_y + Emu(bar_h).inches + 0.35
            legend_h_in = 0.85 + 0.4 * len(status_legend_rows)
            add_ivory_panel(
                status_slide, MARGIN_X, Inches(legend_y_in), CONTENT_W, Inches(legend_h_in),
                "%", status_legend_title, status_legend_rows, mode="legend", theme=ctx.theme,
            )
            status_content_bottom_in = legend_y_in + legend_h_in
    else:
        add_native_bar_chart(
            status_slide, MARGIN_X, Inches(status_body_y), CONTENT_W, status_body_h,
            block["categories"], block["values"], colors=[ctx.accent_bar_color],
        )
    if status_has_caption:
        add_ai_insight_strip(status_slide, MARGIN_X, Inches(status_content_bottom_in + 0.12), CONTENT_W, block["ai_caption"])
    return status_slide


def _build_critical_table_slide(block: dict, ctx: _PptBlockContext):
    table_slide = ctx.prs.slides.add_slide(ctx.prs.slide_layouts[6])
    add_logo(table_slide, ctx.logo_path)
    # RED_CRIT TIDAK ikut tema (severity fixed) — cuma cabang "tidak kritis" yang ikut tema.
    kicker_color = RED_CRIT if block["kicker_is_critical"] else ctx.accent_main
    add_kicker(table_slide, block["kicker"], color=kicker_color)
    title_bottom = add_title(table_slide, block["title"])

    table_y = max(title_bottom + 0.15, 1.55)
    add_native_table(table_slide, MARGIN_X, Inches(table_y), CONTENT_W, Inches(4.9), block["headers"], block["rows"], set(block["highlight_idx"]), theme=ctx.theme)

    if block["caption"]:
        cap_box = table_slide.shapes.add_textbox(MARGIN_X, Inches(table_y) + Inches(4.9) + Inches(0.1), CONTENT_W, Inches(0.4))
        cp2 = cap_box.text_frame.paragraphs[0]
        cp2.text = block["caption"]
        _set_font(cp2, BODY_FONT, Pt(10.5), italic=True, color=GRAY_TEXT)
    return table_slide


def _build_asset_cards_slide(block: dict, ctx: _PptBlockContext):
    asset_slide = ctx.prs.slides.add_slide(ctx.prs.slide_layouts[6])
    use_podium = ctx.asset_style == "podium" and len(block["items"]) == 3
    if use_podium:
        # Titik variasi tampilan: ranking podium top-3 (bg TERANG, analog cover
        # normal) — alternatif dari baris kartu gelap standar, lihat add_podium_row.
        add_logo(asset_slide, ctx.logo_path)
        add_kicker(asset_slide, block["kicker"], color=ctx.accent_main)
        title_bottom = add_title(asset_slide, block["title"], color=TEXT_DARK)
        podium_y = max(title_bottom + 0.3, 2.0)
        add_podium_row(asset_slide, MARGIN_X, Inches(podium_y), CONTENT_W, Inches(max(3.0, 6.5 - podium_y)), block["items"], theme=ctx.theme)
    elif ctx.asset_style == "bars":
        # Titik variasi tampilan: daftar berperingkat dgn batang proporsional (bg
        # GELAP, analog kartu standar) — dipakai utk jumlah item berapa pun (bukan
        # cuma tepat 3 seperti podium), lihat add_asset_ranked_bars.
        add_dark_bg(asset_slide, theme=ctx.theme)
        add_logo(asset_slide, ctx.logo_path)
        add_kicker(asset_slide, block["kicker"], color=ctx.accent_light)
        title_bottom = add_title(asset_slide, block["title"], color=WHITE)
        bars_y = max(title_bottom + 0.3, 2.0)
        # Sisa ruang di bawah baris terakhir dibagi rata (digeser ke tengah), bukan
        # ditumpuk semua di bawah — pola sama dgn add_asset_card_row (cards).
        bars_row_h_in = 1.0
        bars_content_h_in = len(block["items"]) * bars_row_h_in
        bars_available_in = 6.6 - bars_y
        if bars_available_in > bars_content_h_in:
            bars_y += (bars_available_in - bars_content_h_in) / 2
        add_asset_ranked_bars(asset_slide, MARGIN_X, Inches(bars_y), CONTENT_W, block["items"], row_h=Inches(bars_row_h_in), theme=ctx.theme)
    else:
        add_dark_bg(asset_slide, theme=ctx.theme)
        add_logo(asset_slide, ctx.logo_path)
        add_kicker(asset_slide, block["kicker"], color=ctx.accent_light)
        title_bottom = add_title(asset_slide, block["title"], color=WHITE)

        card_items = [(it["num"], it["name"], it["stat"], it["detail"]) for it in block["items"]]
        cards_y = max(title_bottom + 0.2, 1.9)
        add_asset_card_row(asset_slide, MARGIN_X, Inches(cards_y), CONTENT_W, Inches(max(3.0, 6.5 - cards_y)), card_items, theme=ctx.theme)
    return asset_slide


def _build_key_findings_slide(block: dict, ctx: _PptBlockContext):
    find_slide = ctx.prs.slides.add_slide(ctx.prs.slide_layouts[6])
    add_logo(find_slide, ctx.logo_path)
    add_kicker(find_slide, block["kicker"], color=ctx.accent_main)
    title_bottom = add_title(find_slide, block["title"])

    findings_items = [(it["num"], it["title"], it["detail"]) for it in block["items"]]

    # RED_CRIT TIDAK ikut tema (severity fixed) — cuma cabang "tidak kritis" yang ikut tema.
    def _finding_color(idx, item, _items=block["items"]):
        return RED_CRIT if _items[idx]["is_critical"] else ctx.accent_main

    # Keluhan nyata dari pengguna: 1 kolom penuh CONTENT_W bikin baris ~11.8in
    # lebar (nyaris selebar slide) sekaligus font kekecilan (9.3pt) kalau itemnya
    # banyak — grid 2 kolom dipakai begitu item > 2 (lihat catatan di add_badge_list).
    findings_cols = 2 if len(findings_items) > 2 else 1
    add_badge_list(find_slide, MARGIN_X, Inches(max(title_bottom + 0.2, 1.7)), CONTENT_W, findings_items, badge_color=_finding_color, row_h=Inches(1.0), max_y=SLIDE_H - Inches(0.4), cols=findings_cols)
    return find_slide


def _build_recommendations_slide(block: dict, ctx: _PptBlockContext):
    rec_slide = ctx.prs.slides.add_slide(ctx.prs.slide_layouts[6])
    add_logo(rec_slide, ctx.logo_path)
    add_kicker(rec_slide, block["kicker"], color=ctx.accent_main)
    title_bottom = add_title(rec_slide, block["title"])

    items = block["items"]
    start_y_rec = Inches(max(title_bottom + 0.2, 1.7))

    if ctx.recommendation_style == "timeline" and 2 <= len(items) <= 6:
        # Titik variasi tampilan: garis waktu horizontal (analog timeline di
        # export_pdf.py) — alternatif dari grid kartu standar di bawah, dipakai
        # kalau jumlah rekomendasi pas 2-6 (lihat add_recommendation_timeline).
        timeline_h = Inches(max(2.6, 6.9 - Emu(start_y_rec).inches))
        add_recommendation_timeline(rec_slide, MARGIN_X, start_y_rec, CONTENT_W, timeline_h, items, theme=ctx.theme)
        return rec_slide

    if ctx.recommendation_style == "banners":
        # Titik variasi tampilan: daftar banner selebar halaman bertumpuk vertikal
        # (analog _recommendation_banner_list_html di export_pdf.py) — alternatif
        # dari grid kartu, dipakai utk jumlah item BERAPA PUN (beda dari timeline
        # yang dibatasi 2-6).
        add_recommendation_banner_list(rec_slide, MARGIN_X, start_y_rec, CONTENT_W, items, max_y=SLIDE_H - Inches(0.4), theme=ctx.theme)
        return rec_slide

    card_cols = ctx.card_cols
    gap = Inches(0.3)
    rec_rows = math.ceil(len(items) / card_cols)

    # Pre-pass: hitung tinggi SETIAP baris (di skala normal) sebelum menggambar
    # apapun, supaya total tinggi semua baris bisa diketahui DULU. BUG YANG
    # DIPERBAIKI: sebelumnya tidak ada pengecekan terhadap SLIDE_H — laporan
    # dengan 3 baris (6 rekomendasi, 2 kolom) yang detailnya panjang bisa membuat
    # baris terakhir meluber ke luar slide, persis pola yang sama dengan bug
    # Temuan Utama. Kalau totalnya bakal melebihi ruang tersisa, seluruh baris
    # (font & tinggi kartu) dikecilkan proporsional supaya semua rekomendasi
    # tetap tampil penuh di dalam slide.
    def _row_items_and_width(r):
        row_start = r * card_cols
        row_items = items[row_start:row_start + card_cols]
        row_items_count = len(row_items)
        card_w = (CONTENT_W - gap * (row_items_count - 1)) / row_items_count
        return row_items, card_w

    def _estimate_row_height_in(row_items, card_w, title_pt, detail_pt):
        text_w_in = Emu(card_w - Inches(0.4)).inches
        row_height_in = 1.3
        for item in row_items:
            title_h_in = _estimate_wrapped_height_in(item["title"], title_pt, text_w_in)
            content_h_in = 0.62 + title_h_in + 0.25
            if item["detail"]:
                detail_h_in = _estimate_wrapped_height_in(item["detail"], detail_pt, text_w_in)
                content_h_in += detail_h_in + 0.1
            row_height_in = max(row_height_in, content_h_in)
        return row_height_in

    available_in = Emu(SLIDE_H - start_y_rec).inches - 0.4
    est_total_in = sum(
        _estimate_row_height_in(_row_items_and_width(r)[0], _row_items_and_width(r)[1], 13, 10.5) + Emu(gap).inches
        for r in range(rec_rows)
    ) - Emu(gap).inches
    scale = 1.0
    if available_in > 0 and est_total_in > available_in:
        scale = max(available_in / est_total_in, 0.55)
    title_pt, detail_pt = 13 * scale, 10.5 * scale

    cur_y_rec = start_y_rec
    # Keluhan nyata dari pengguna: kalau rekomendasinya sedikit/singkat (est_total_in
    # jauh di bawah available_in, scale tetap 1.0), grid kartu cuma menempati bagian
    # atas slide, sisanya kosong sampai footer. Geser TITIK AWAL baris pertama ke
    # bawah supaya sisa ruang kosong terbagi rata atas & bawah alih-alih menumpuk di
    # bawah — sama seperti perbaikan pada slide Ringkasan Eksekutif & Sorotan Data.
    if available_in > 0 and est_total_in < available_in:
        cur_y_rec += Inches((available_in - est_total_in) / 2)
    for r in range(rec_rows):
        row_items, card_w = _row_items_and_width(r)
        row_height_in = _estimate_row_height_in(row_items, card_w, title_pt, detail_pt)
        card_h = Inches(row_height_in)
        for c, item in enumerate(row_items):
            cx = MARGIN_X + c * (card_w + gap)
            cy = cur_y_rec
            card = rec_slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, card_w, card_h)
            card.fill.solid()
            card.fill.fore_color.rgb = IVORY
            card.line.color.rgb = PANEL_BORDER
            card.line.width = Pt(0.75)
            _no_shadow(card)
            _modest_corner(card)
            add_badge_circle(rec_slide, cx + Inches(0.2), cy + Inches(0.18 * scale), Inches(0.36 * scale), item["num"], ctx.accent_light, font_size=Pt(13 * scale))
            title_box = rec_slide.shapes.add_textbox(cx + Inches(0.2), cy + Inches(0.62 * scale), card_w - Inches(0.4), Inches(0.35))
            ttf2 = title_box.text_frame
            ttf2.word_wrap = True
            tp2 = ttf2.paragraphs[0]
            tp2.text = item["title"]
            _set_font(tp2, BODY_FONT, Pt(title_pt), bold=True, color=TEXT_DARK)
            if item["detail"]:
                det_box = rec_slide.shapes.add_textbox(cx + Inches(0.2), cy + Inches(0.95 * scale), card_w - Inches(0.4), card_h - Inches(1.0 * scale))
                dtf2 = det_box.text_frame
                dtf2.word_wrap = True
                dp2 = dtf2.paragraphs[0]
                dp2.text = item["detail"]
                _set_font(dp2, BODY_FONT, Pt(detail_pt), color=GRAY_TEXT)
        cur_y_rec += card_h + gap
    return rec_slide


def _build_conclusion_slide(block: dict, ctx: _PptBlockContext):
    concl_slide = ctx.prs.slides.add_slide(ctx.prs.slide_layouts[6])
    add_dark_bg(concl_slide, theme=ctx.theme)
    add_logo(concl_slide, ctx.logo_path)
    add_kicker(concl_slide, block["kicker"], color=ctx.accent_light)
    title_bottom = add_title(concl_slide, block["title"], color=WHITE)

    left_w2 = Inches(7.0)
    content_top2 = max(title_bottom + 0.15, 1.7)
    priority_items = [(p["letter"], p["text"]) for p in block["priority_items"]]
    # panel_side cuma dipakai kalau ada priority_items — kalau kosong, swap ke
    # "left" akan menyisakan kolom kiri kosong & teks malah ke kanan (lebih buruk
    # dari layout defaultnya), sama seperti guard yang sama di export_pdf.py.
    if priority_items and ctx.panel_side == "left":
        text_x2 = MARGIN_X + Inches(4.3) + Inches(0.4)
        panel_x3 = MARGIN_X
    else:
        text_x2 = MARGIN_X
        panel_x3 = MARGIN_X + left_w2 + Inches(0.4)
    para_box2 = concl_slide.shapes.add_textbox(text_x2, Inches(content_top2), left_w2, Inches(1.8))
    ptf2 = para_box2.text_frame
    ptf2.word_wrap = True
    pp2 = ptf2.paragraphs[0]
    pp2.text = block["text"]
    _set_font(pp2, BODY_FONT, Pt(13), color=RGBColor(0xE8, 0xEC, 0xE6))

    pill_y = Inches(content_top2 + 2.0)
    for pill_text in block["pills"][:3]:
        add_pill_stat(concl_slide, text_x2, pill_y, left_w2, Inches(0.6), pill_text, theme=ctx.theme)
        pill_y += Inches(0.75)

    if priority_items:
        add_priority_panel(
            concl_slide, panel_x3, Inches(content_top2), Inches(4.3), Inches(max(3.0, 6.3 - content_top2)),
            block["priority_panel_title"], priority_items, theme=ctx.theme,
        )
    return concl_slide


def _build_closing_slide(block: dict, ctx: _PptBlockContext):
    if ctx.cover_style == "split":
        closing_block = {**block, "hero_stat": ctx.cover_hero_stat}
        add_split_closing_slide(ctx.prs, closing_block, ctx.flourish_corner, theme=ctx.theme)
        return None
    closing = ctx.prs.slides.add_slide(ctx.prs.slide_layouts[6])
    add_dark_bg(closing, theme=ctx.theme)
    add_corner_flourish(closing, ctx.flourish_corner, theme=ctx.theme)
    # Teks digeser ke kanan kalau flourish-nya di pojok kiri-bawah (satu-satunya
    # varian yang jangkauannya menjorok ke area teks, yang start dari MARGIN_X) —
    # keluhan nyata dari pengguna: garis lengkung dekoratif menembus teks penutup.
    text_x2 = Inches(2.3) if ctx.flourish_corner == "bottom_left" else MARGIN_X
    text_w2 = Inches(9) - (text_x2 - MARGIN_X)
    text_w2_in = Emu(text_w2).inches
    title_box2 = closing.shapes.add_textbox(text_x2, Inches(3.0), text_w2, Inches(1.0))
    tp3 = title_box2.text_frame.paragraphs[0]
    tp3.text = block["thank_you"]
    _set_font(tp3, TITLE_FONT, Pt(40), bold=True, color=WHITE)
    sub2_top_in = 3.85
    # BUG YANG DIPERBAIKI (dilaporkan user): box ini dulu tinggi TETAP 0.8in,
    # sementara note_box di bawah diposisikan berdasarkan estimasi tinggi KONTEN
    # sebenarnya (sub2_height_in, sering < 0.8in utk judul 1 baris) — box tetap
    # 0.8in penuh jadi meluber menimpa note_box. Tinggi box sekarang eksplisit
    # mengikuti estimasi yang sama dipakai utk memposisikan note_box di bawahnya.
    sub2_height_in = _estimate_wrapped_height_in(block["title"], 14, text_w2_in)
    sub_box2 = closing.shapes.add_textbox(text_x2, Inches(sub2_top_in), text_w2, Inches(sub2_height_in + 0.1))
    sp2 = sub_box2.text_frame
    sp2.word_wrap = True
    sp2_p = sp2.paragraphs[0]
    sp2_p.text = block["title"]
    _set_font(sp2_p, BODY_FONT, Pt(14), color=WHITE)
    note_top_in = max(sub2_top_in + sub2_height_in + 0.1, 4.4)
    note_box = closing.shapes.add_textbox(text_x2, Inches(note_top_in), text_w2, Inches(0.4))
    np_ = note_box.text_frame.paragraphs[0]
    np_.text = block["note"]
    _set_font(np_, BODY_FONT, Pt(11.5), italic=True, color=ctx.accent_soft)
    return None


_PPT_BLOCK_BUILDERS = {
    "cover": _build_cover_slide,
    "intro": _build_intro_slide,
    "executive_summary": _build_executive_summary_slide,
    "dynamic_section": _build_dynamic_section_slide,
    "category_distribution": _build_category_distribution_slide,
    "severity_distribution": _build_severity_distribution_slide,
    "status_distribution": _build_status_distribution_slide,
    "critical_table": _build_critical_table_slide,
    "asset_cards": _build_asset_cards_slide,
    "key_findings": _build_key_findings_slide,
    "recommendations": _build_recommendations_slide,
    "conclusion": _build_conclusion_slide,
    "closing": _build_closing_slide,
}


class PPTXExporter:
    @classmethod
    def generate_ppt_report(cls, report: Report) -> bytes:
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H

        logo_path = _resolve_logo_path()
        blocks = build_report_blocks(report)

        # Varian tampilan (cover_style, category_style, dst) DIBACA dari report.visual_style,
        # BUKAN di-random di sini lagi — BUG YANG DIPERBAIKI (dilaporkan user): dulu tiap kali
        # generate_ppt_report dipanggil, pilihan acak baru diambil, jadi preview web (yang
        # membaca kombinasi TETAP per laporan lewat endpoint /blocks) bisa menampilkan bentuk
        # yang beda dari file yang benar-benar diunduh. Sekarang preview & export SAMA-SAMA
        # baca report.visual_style yang SUDAH DIKUNCI sekali sewaktu analisis AI berhasil (lihat
        # pick_visual_style() di report_render_logic.py) — dijamin identik utk 1 laporan yang
        # sama, dan React (ReportBlockRenderer.tsx) sekarang genuinely merender SEMUA varian ini
        # (bukan lagi cuma 1 bentuk tetap), jadi preview akurat mencerminkan file yang diunduh.
        vs = get_visual_style(report)
        flourish_corner = vs["flourish_corner"]
        panel_side = vs["panel_side"]
        stat_cols = vs["stat_cols"]
        card_cols = vs["card_cols"]
        category_style = vs["category_style"]
        status_style = vs["status_style"]
        cover_style = vs["cover_style"]
        asset_style = vs["asset_style"]
        recommendation_style = vs["recommendation_style"]
        # Teks kicker TETAP (bukan bagian dari visual_style — cuma variasi kata, bukan bentuk)
        kicker_ringkasan = "Executive Summary" if is_english(report) else "Ringkasan Eksekutif"
        kicker_analisis = "DATA ANALYSIS" if is_english(report) else "ANALISIS DATA"

        # Palet warna tema (report.theme_color) — accent_bar_color TIDAK LAGI dipakai untuk
        # warna aksen (dulu diacak hijau/emas lewat visual_style, independen dari pilihan user
        # di Report Settings). Sekarang accent_bar_color = ctx.accent_main, konsisten dgn tema.
        theme_key = resolve_theme_color(report)
        palette = THEME_PALETTES[theme_key]
        accent_bar_color = palette["main"]

        content_slides: list = []  # dipakai utk stamping footer di akhir (kecuali cover/penutup)

        ctx = _PptBlockContext(
            report=report,
            prs=prs,
            logo_path=logo_path,
            panel_side=panel_side,
            stat_cols=stat_cols,
            card_cols=card_cols,
            flourish_corner=flourish_corner,
            accent_bar_color=accent_bar_color,
            category_style=category_style,
            status_style=status_style,
            cover_style=cover_style,
            asset_style=asset_style,
            recommendation_style=recommendation_style,
            accent_main=palette["main"],
            accent_bg=palette["bg"],
            accent_chart=palette["chart"],
            accent_light=palette["light"],
            accent_soft=palette["soft"],
            theme=palette,
            kicker_ringkasan=kicker_ringkasan,
            kicker_analisis=kicker_analisis,
        )

        for block in blocks:
            builder = _PPT_BLOCK_BUILDERS.get(block["kind"])
            if builder:
                slide = builder(block, ctx)
                if slide is not None:
                    content_slides.append(slide)

        # -------------------------------------------------------------
        # Footer: nomor halaman, semua slide isi kecuali cover & penutup
        # -------------------------------------------------------------
        total_pages = len(content_slides)
        for idx, s in enumerate(content_slides):
            add_footer(s, idx + 1, total_pages)

        ppt_io = __import__("io").BytesIO()
        prs.save(ppt_io)
        ppt_io.seek(0)
        return ppt_io.read()


