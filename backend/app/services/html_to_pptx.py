"""
Konverter HTML → python-pptx.

python-pptx TIDAK bisa membaca HTML sama sekali — dia cuma punya API level rendah
(paragraph, run, table shape). Modul ini menjembatani itu: mem-parse HTML sederhana
yang dihasilkan rich text editor internal (Bold/Italic/Underline/warna/highlight/align/
list/tabel/link) menjadi struktur "blocks" yang lalu digambar manual ke slide PPTX.

CATATAN KETERBATASAN (biar jelas, bukan janji "identik 100% dengan tampilan web"):
- Highlight warna (mark/background-color) TIDAK didukung python-pptx untuk teks biasa,
  jadi diabaikan saat render ke PPTX (tetap muncul normal di PDF & preview web).
- Tabel di dalam teks narasi dirender sebagai shape tabel terpisah di bawah teks pada
  slide yang sama, bukan "inline" persis di tengah paragraf (PowerPoint memang tidak
  punya konsep tabel inline di dalam text box seperti web/PDF).
- Parser ini sengaja sederhana (bukan HTML parser umum/lengkap) karena sumber HTML-nya
  selalu dari editor internal sendiri (TipTap), bukan HTML sembarangan dari luar.
"""

from html.parser import HTMLParser
from html import unescape
import re

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


_COLOR_HEX_RE = re.compile(r"#([0-9a-fA-F]{6})")


def _parse_color_from_style(style: str):
    """Ambil warna teks (color: #rrggbb) dari atribut style inline, kalau ada."""
    if not style:
        return None
    m = re.search(r"color\s*:\s*(#[0-9a-fA-F]{6})", style)
    if not m:
        return None
    hex_match = _COLOR_HEX_RE.search(m.group(1))
    if not hex_match:
        return None
    return RGBColor.from_string(hex_match.group(1))


class _HtmlToBlocks(HTMLParser):
    """
    Parser satu-arah: menghasilkan self.blocks, list dari salah satu bentuk berikut:
      {"type": "paragraph", "align": "left"|"center"|"right", "list": None|"bullet"|"number",
       "runs": [{"text": str, "bold": bool, "italic": bool, "underline": bool,
                 "color": RGBColor|None, "link": str|None}]}
      {"type": "table", "rows": [[cell_text, ...], ...]}
    """

    _INLINE_TAGS = {"strong", "b", "em", "i", "u", "mark", "a", "span"}
    _BLOCK_TAGS = {"p", "div", "li", "h1", "h2", "h3"}

    def __init__(self):
        super().__init__()
        self.blocks = []
        self._style_stack = []  # tumpukan dict(bold, italic, underline, color, link)
        self._current_runs = []
        self._current_align = "left"
        self._list_stack = []  # tumpukan 'bullet' / 'number'
        self._in_table = False
        self._table_rows = []
        self._current_row = []
        self._current_cell_text = []
        self._in_cell = False

    # ── Helpers ──────────────────────────────────────────────────────────
    def _current_style(self):
        merged = {"bold": False, "italic": False, "underline": False, "color": None, "link": None}
        for s in self._style_stack:
            merged.update({k: v for k, v in s.items() if v})
        return merged

    def _flush_paragraph(self, list_type=None):
        has_content = any(r["text"].strip() for r in self._current_runs)
        if has_content:
            self.blocks.append({
                "type": "paragraph",
                "align": self._current_align,
                "list": list_type,
                "runs": self._current_runs,
            })
        self._current_runs = []
        self._current_align = "left"

    # ── HTMLParser overrides ─────────────────────────────────────────────
    def handle_starttag(self, tag, attrs):
        attrs_dict = dict(attrs)
        if self._in_table:
            if tag == "tr":
                self._current_row = []
            elif tag in ("td", "th"):
                self._in_cell = True
                self._current_cell_text = []
            return

        if tag == "table":
            self._flush_paragraph()
            self._in_table = True
            self._table_rows = []
            return

        if tag in ("ul", "ol"):
            self._list_stack.append("bullet" if tag == "ul" else "number")
            return

        if tag == "li":
            self._flush_paragraph()
            return

        if tag == "br":
            self._current_runs.append({"text": "\n", "bold": False, "italic": False, "underline": False, "color": None, "link": None})
            return

        if tag in self._BLOCK_TAGS:
            style = attrs_dict.get("style", "")
            if "text-align: center" in style or "text-align:center" in style:
                self._current_align = "center"
            elif "text-align: right" in style or "text-align:right" in style:
                self._current_align = "right"
            elif "text-align: justify" in style or "text-align:justify" in style:
                self._current_align = "justify"
            return

        if tag in self._INLINE_TAGS:
            style_attrs = {"bold": False, "italic": False, "underline": False, "color": None, "link": None}
            if tag in ("strong", "b"):
                style_attrs["bold"] = True
            elif tag in ("em", "i"):
                style_attrs["italic"] = True
            elif tag == "u":
                style_attrs["underline"] = True
            elif tag == "a":
                style_attrs["underline"] = True
                style_attrs["link"] = attrs_dict.get("href")
            elif tag == "span":
                color = _parse_color_from_style(attrs_dict.get("style", ""))
                if color:
                    style_attrs["color"] = color
            self._style_stack.append(style_attrs)

    def handle_endtag(self, tag):
        if self._in_table:
            if tag in ("td", "th"):
                self._in_cell = False
                self._current_row.append("".join(self._current_cell_text).strip())
            elif tag == "tr":
                if self._current_row:
                    self._table_rows.append(self._current_row)
            elif tag == "table":
                self._in_table = False
                if self._table_rows:
                    self.blocks.append({"type": "table", "rows": self._table_rows})
            return

        if tag in ("ul", "ol"):
            if self._list_stack:
                self._list_stack.pop()
            return

        if tag == "li":
            list_type = self._list_stack[-1] if self._list_stack else "bullet"
            self._flush_paragraph(list_type=list_type)
            return

        if tag in self._BLOCK_TAGS:
            self._flush_paragraph()
            return

        if tag in self._INLINE_TAGS:
            if self._style_stack:
                self._style_stack.pop()

    def handle_data(self, data):
        if not data:
            return
        if self._in_table:
            if self._in_cell:
                self._current_cell_text.append(data)
            return
        # Buang whitespace di awal akumulasi paragraf baru (misal newline/indentasi
        # yang kebetulan ada di antara tag), tapi tetap simpan spasi di tengah kalimat.
        if not self._current_runs and not data.strip():
            return
        style = self._current_style()
        self._current_runs.append({
            "text": data,
            "bold": style["bold"],
            "italic": style["italic"],
            "underline": style["underline"],
            "color": style["color"],
            "link": style["link"],
        })

    def close(self):
        super().close()
        self._flush_paragraph()


def parse_html_to_blocks(html: str):
    """Entry point: HTML string -> list of block dicts (lihat docstring _HtmlToBlocks)."""
    if not html:
        return []
    parser = _HtmlToBlocks()
    parser.feed(unescape(html))
    parser.close()
    return parser.blocks


_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}


def render_blocks_to_textframe(text_frame, blocks, base_size=Pt(14), base_color=None, list_indent=True):
    """
    Menggambar blocks (kecuali tabel) ke sebuah text_frame python-pptx yang sudah ada.
    Mengembalikan list tabel (block bertipe "table") yang TIDAK ikut digambar di sini,
    supaya caller bisa merender tabel itu sebagai shape terpisah di slide yang sama.
    """
    pending_tables = []
    first_paragraph = True
    bullet_counters = {}

    for block in blocks:
        if block["type"] == "table":
            pending_tables.append(block)
            continue

        p = text_frame.paragraphs[0] if first_paragraph else text_frame.add_paragraph()
        first_paragraph = False
        p.alignment = _ALIGN_MAP.get(block["align"], PP_ALIGN.LEFT)
        p.space_after = Pt(10)
        p.line_spacing = 1.25

        prefix = ""
        if block["list"] == "bullet" and list_indent:
            prefix = "•  "
        elif block["list"] == "number" and list_indent:
            bullet_counters[id(blocks)] = bullet_counters.get(id(blocks), 0) + 1
            prefix = f"{bullet_counters[id(blocks)]}.  "

        runs = block["runs"] or [{"text": "", "bold": False, "italic": False, "underline": False, "color": None, "link": None}]

        if prefix:
            r0 = p.add_run()
            r0.text = prefix
            r0.font.size = base_size
            r0.font.bold = True
            if base_color:
                r0.font.color.rgb = base_color

        for seg in runs:
            if seg["text"] == "":
                continue
            # "\n" ditulis sebagai paragraf baru (soft break) supaya kompatibel di semua versi python-pptx
            parts = seg["text"].split("\n")
            for idx, part in enumerate(parts):
                if idx > 0:
                    p = text_frame.add_paragraph()
                    p.alignment = _ALIGN_MAP.get(block["align"], PP_ALIGN.LEFT)
                if part == "":
                    continue
                run = p.add_run()
                run.text = part
                run.font.size = base_size
                run.font.bold = seg["bold"]
                run.font.italic = seg["italic"]
                run.font.underline = seg["underline"]
                if seg["color"]:
                    run.font.color.rgb = seg["color"]
                elif base_color:
                    run.font.color.rgb = base_color
                if seg["link"]:
                    try:
                        run.hyperlink.address = seg["link"]
                    except Exception:
                        pass

    if first_paragraph:
        # Tidak ada konten sama sekali — biarkan paragraf pertama kosong tapi valid
        text_frame.paragraphs[0].text = ""

    return pending_tables


def render_tables_to_slide(slide, tables, left, top, width, row_height=Inches(0.35)):
    """
    Menggambar tabel-tabel (hasil dari render_blocks_to_textframe) sebagai shape tabel
    python-pptx yang sesungguhnya, disusun vertikal ke bawah mulai dari posisi `top`.
    Mengembalikan posisi Y akhir (Emu) setelah tabel terakhir, untuk elemen berikutnya.
    """
    current_top = top
    for table_block in tables:
        rows_data = table_block["rows"]
        if not rows_data:
            continue
        n_rows = len(rows_data)
        n_cols = max(len(r) for r in rows_data)
        height = row_height * n_rows

        graphic_frame = slide.shapes.add_table(n_rows, n_cols, left, current_top, width, height)
        table = graphic_frame.table

        for r_idx, row_data in enumerate(rows_data):
            for c_idx in range(n_cols):
                cell_text = row_data[c_idx] if c_idx < len(row_data) else ""
                cell = table.cell(r_idx, c_idx)
                cell.text = cell_text
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(10)
                    if r_idx == 0:
                        paragraph.font.bold = True

        current_top = current_top + height + Inches(0.25)

    return current_top